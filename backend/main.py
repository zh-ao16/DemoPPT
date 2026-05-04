#!/usr/bin/env python3
"""DemoPPT 后端API服务 v2.0 - 专业级PPT生成"""
from fastapi import FastAPI, HTTPException, Request, Header, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional, List
import asyncio
import os
from pathlib import Path
import uuid
import json
import re
from knowledge_base import kb, parse_document, KnowledgeBase
from industry_data import INDUSTRY_KB, get_industry_context, search_and_get_context
from color_palette import generate_palette, get_preset_palette, list_preset_palettes

# 进度状态存储（内存字典，key为session_id）
progress_store = {}

app = FastAPI(title="DemoPPT API", version="2.3.0")

# CORS配置 - 生产环境限制来源
ALLOWED_ORIGINS = os.environ.get("ALLOWED_ORIGINS", "http://localhost:3000,http://localhost:5173").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

# 目录
BASE_DIR = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# 请求模型
class OutlineRequest(BaseModel):
    topic: str
    requirements: Optional[str] = ""
    industry: Optional[str] = "general"
    model_config_id: Optional[int] = None  # 用户配置的模型ID

class ContentRequest(BaseModel):
    topic: str
    outline: List[dict]
    template: str = "academic"
    industry: Optional[str] = "general"
    # 演讲者备注
    speaker_notes: Optional[bool] = False
    # 品牌定制
    brand_name: Optional[str] = ""
    brand_logo: Optional[str] = ""  # logo URL或base64
    brand_color: Optional[str] = ""  # 自定义主题色
    # 多语言
    language: Optional[str] = "zh"  # zh/en/zh-TW
    # AI模型配置（用户自定义）
    model_config_id: Optional[int] = None  # 用户配置的模型ID

class DigitalHumanRequest(BaseModel):
    text: str
    avatar: str = "default"

class DocumentConvertRequest(BaseModel):
    text: str
    industry: str = "general"
    template: str = "business"
    model_config_id: Optional[int] = None  # 用户配置的模型ID

# 行业模板配置
INDUSTRY_TEMPLATES = {
    "education": {
        "name": "教育培训",
        "templates": ["academic", "nature", "sky"],
        "color": "#276749"
    },
    "medical": {
        "name": "医疗健康",
        "templates": ["simple", "nature", "ocean"],
        "color": "#0077b6"
    },
    "ecommerce": {
        "name": "电商零售",
        "templates": ["gradient", "festive", "cyber"],
        "color": "#ff6b6b"
    },
    "finance": {
        "name": "金融投资",
        "templates": ["business", "navy_gold", "classic_blue"],
        "color": "#2c5282"
    },
    "technology": {
        "name": "科技互联网",
        "templates": ["tech", "cyber", "future"],
        "color": "#1a1a4e"
    },
    "government": {
        "name": "政府企业",
        "templates": ["enterprise", "royal", "academic"],
        "color": "#4c1d95"
    },
    "realestate": {
        "name": "房产建筑",
        "templates": ["business", "gray_elegant", "ink_wash"],
        "color": "#553c9a"
    },
    "media": {
        "name": "传媒广告",
        "templates": ["gradient", "cyber", "royal"],
        "color": "#ff00ff"
    },
    "manufacture": {
        "name": "制造业",
        "templates": ["enterprise", "nature", "classic_blue"],
        "color": "#3182ce"
    },
    "general": {
        "name": "通用场景",
        "templates": ["business", "simple", "elegant"],
        "color": "#2d3748"
    }
}

# AI模型配置
AI_MODELS = {
    "deepseek": {
        "name": "DeepSeek",
        "provider": "深度求索",
        "api_url": "https://api.deepseek.com/v1/chat/completions",
        "models": ["deepseek-chat", "deepseek-coder"],
        "default_model": "deepseek-chat",
        "env_key": "DEEPSEEK_API_KEY",
        "supports_functions": True,
        "max_tokens": 16384,
    },
    "doubao": {
        "name": "豆包",
        "provider": "字节跳动",
        "api_url": "https://ark.cn-beijing.volces.com/api/v3/chat/completions",
        "models": ["doubao-pro-32k", "doubao-pro-128k"],
        "default_model": "doubao-pro-32k",
        "env_key": "DOUBAO_API_KEY",
        "supports_functions": False,
        "max_tokens": 32000,
    },
    "zhipu": {
        "name": "智谱GLM",
        "provider": "智谱AI",
        "api_url": "https://open.bigmodel.cn/api/paas/v4/chat/completions",
        "models": ["glm-4", "glm-4-flash", "glm-3-turbo"],
        "default_model": "glm-4",
        "env_key": "ZHIPU_API_KEY",
        "supports_functions": True,
        "max_tokens": 128000,
    },
    "qwen": {
        "name": "通义千问",
        "provider": "阿里云",
        "api_url": "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions",
        "models": ["qwen-turbo", "qwen-plus", "qwen-max"],
        "default_model": "qwen-turbo",
        "env_key": "DASHSCOPE_API_KEY",
        "supports_functions": True,
        "max_tokens": 8192,
    },
    "wenxin": {
        "name": "文心一言",
        "provider": "百度",
        "api_url": "https://qianfan.baidubce.com/v2/chat/completions",
        "models": ["ernie-4.0-8k-latest", "ernie-3.5-8k-latest"],
        "default_model": "ernie-3.5-8k-latest",
        "env_key": "WENXIN_API_KEY",
        "supports_functions": True,
        "max_tokens": 8192,
    },
    "hunyuan": {
        "name": "腾讯混元",
        "provider": "腾讯云",
        "api_url": "https://hunyuan.cloud.tencent.com/v1/chat/completions",
        "models": ["hunyuan-pro", "hunyuan-standard"],
        "default_model": "hunyuan-standard",
        "env_key": "HUNYUAN_API_KEY",
        "supports_functions": False,
        "max_tokens": 4096,
    },
    "spark": {
        "name": "讯飞星火",
        "provider": "科大讯飞",
        "api_url": "https://spark-api.xf-yun.com/v3.5/chat",
        "models": ["v3.5", "v3.0", "v2.0"],
        "default_model": "v3.5",
        "env_key": "SPARK_API_KEY",
        "supports_functions": False,
        "max_tokens": 8192,
    },
    "minimax": {
        "name": "MiniMax",
        "provider": "MiniMax",
        "api_url": "https://api.minimax.chat/v1/text/chatcompletion_v2",
        "models": ["abab6-chat"],
        "default_model": "abab6-chat",
        "env_key": "MINIMAX_API_KEY",
        "supports_functions": False,
        "max_tokens": 8192,
    },
    "openai": {
        "name": "OpenAI GPT",
        "provider": "OpenAI",
        "api_url": "https://api.openai.com/v1/chat/completions",
        "models": ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-3.5-turbo"],
        "default_model": "gpt-4o-mini",
        "env_key": "OPENAI_API_KEY",
        "supports_functions": True,
        "max_tokens": 128000,
    },
    "claude": {
        "name": "Claude",
        "provider": "Anthropic",
        "api_url": "https://api.anthropic.com/v1/messages",
        "models": ["claude-sonnet-4-6", "claude-3-5-sonnet-latest", "claude-3-opus-latest"],
        "default_model": "claude-3-5-sonnet-latest",
        "env_key": "ANTHROPIC_API_KEY",
        "supports_functions": False,
        "max_tokens": 200000,
    },
    "gemini": {
        "name": "Gemini",
        "provider": "Google",
        "api_url": "https://generativelanguage.googleapis.com/v1beta/models",
        "models": ["gemini-1.5-pro", "gemini-1.5-flash", "gemini-1.0-pro"],
        "default_model": "gemini-1.5-flash",
        "env_key": "GOOGLE_API_KEY",
        "supports_functions": True,
        "max_tokens": 1000000,
    },
}

DEFAULT_MODEL = "deepseek"

MODEL_API_FORMATS = {
    "openai": {"format": "openai", "auth_header": "Bearer"},
    "claude": {"format": "claude", "auth_header": "Bearer"},
    "deepseek": {"format": "openai", "auth_header": "Bearer"},
    "doubao": {"format": "openai", "auth_header": "Bearer"},
    "zhipu": {"format": "openai", "auth_header": "Bearer"},
    "qwen": {"format": "openai", "auth_header": "Bearer"},
    "wenxin": {"format": "openai", "auth_header": "Bearer"},
    "hunyuan": {"format": "openai", "auth_header": "Bearer"},
    "minimax": {"format": "openai", "auth_header": "Bearer"},
    "gemini": {"format": "gemini", "auth_header": "Bearer"},
    "spark": {"format": "spark", "auth_header": ""},
}

def call_ai_model(prompt: str, system: str = "你是一个专业的PPT助手", model: str = None, user_api_config: dict = None) -> str:
    """调用AI大模型（支持国内外多种模型 + 用户自定义配置）"""
    import requests

    # 如果用户提供了自己的API配置，优先使用
    if user_api_config:
        try:
            import httpx
            api_base = user_api_config["api_base"].rstrip("/")
            api_key = user_api_config["api_key"]
            model_name = user_api_config["model_name"]
            provider = user_api_config.get("provider", "custom")

            # 跳过明显无效的key
            if not api_key or len(api_key) < 10 or api_key in ("sk-xxxx", "sk-***", ""):
                print(f"用户API Key无效，跳过: {api_key}")
            else:
                headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
                payload = {"model": model_name, "messages": [{"role": "system", "content": system}, {"role": "user", "content": prompt}], "max_tokens": 2000}

                if provider == "anthropic" or "anthropic" in api_base:
                    payload.pop("model")
                    payload["model"] = model_name
                    # DeepSeek Anthropic格式：system必须放顶层，不能在messages里
                    system_msg = None
                    msgs = payload.get("messages", [])
                    if msgs and msgs[0].get("role") == "system":
                        system_msg = msgs[0]["content"]
                        msgs = msgs[1:]
                    payload["messages"] = msgs
                    if system_msg:
                        payload["system"] = system_msg
                    headers["x-api-key"] = api_key
                    headers["anthropic-version"] = "2023-06-01"
                    url = f"{api_base}/messages"
                else:
                    url = f"{api_base}/chat/completions"

                resp = requests.post(url, headers=headers, json=payload, timeout=15)
                resp.raise_for_status()
                result = resp.json()

                if "choices" in result:
                    return result["choices"][0]["message"]["content"]
                elif "content" in result:
                    # content可能是多个block，优先取type=text的
                    for block in result["content"]:
                        if block.get("type") == "text":
                            return block.get("text", "")
                    # 如果没有text block，取第一个block的值
                    return result["content"][0].get("thinking") or result["content"][0].get("text", "")
        except Exception as e:
            print(f"用户API调用失败: {e}，尝试使用系统API")
            # 用户API失败时继续尝试系统API

    # 系统默认API逻辑
    active_model = model or os.environ.get("ACTIVE_MODEL", DEFAULT_MODEL)
    model_config = AI_MODELS.get(active_model, AI_MODELS[DEFAULT_MODEL])

    env_key = model_config["env_key"]
    api_key = os.environ.get(env_key, "")

    if not api_key:
        if active_model == "deepseek":
            api_key = os.environ.get("DEEPSEEK_API_KEY", "")
        elif active_model == "doubao":
            api_key = os.environ.get("DOUBAO_API_KEY", "")

    if not api_key:
        for name, config in AI_MODELS.items():
            key = os.environ.get(config["env_key"], "")
            if key:
                active_model = name
                model_config = config
                api_key = key
                break

    if not api_key:
        print(f"警告：未配置任何AI API Key，使用高质量降级内容")
        return get_fallback_content(prompt, system)

def get_fallback_content(prompt: str, system: str = "") -> str:
    """当无API Key时，返回高质量的降级内容"""
    try:
        from fallback_content import FALLBACK_RESPONSES
    except ImportError:
        return get_generic_fallback(prompt)
    # 根据system和prompt关键词匹配最佳降级内容
    for keywords, content in FALLBACK_RESPONSES:
        if any(kw in system or kw in prompt for kw in keywords):
            return content
    return list(FALLBACK_RESPONSES)[0][1]  # 返回general默认

def get_fallback_outline(topic: str, industry: str) -> list:
    """当无API Key时，返回高质量的大纲"""
    from fallback_content import FALLBACK_OUTLINES
    return FALLBACK_OUTLINES.get(industry, FALLBACK_OUTLINES["general"])

def get_generic_fallback(prompt: str) -> str:
    """通用的降级内容"""
    return """
{
    "title": "核心主题",
    "subtitle": "洞悉本质，引领变革",
    "points": [
        {"heading": "🔍 现状分析", "content": "根据最新行业数据显示，当前市场正经历深刻变革。领先企业正在积极探索新的增长路径，而大多数企业仍在寻找突破方向。关键成功因素包括：战略清晰、执行有力、创新持续、数据驱动。"},
        {"heading": "💡 核心洞察", "content": "研究表明，实现可持续增长的企业通常具备三个共同特征：一是有明确的价值主张，二是有差异化的竞争策略，三是有高效的运营体系。这三者缺一不可，共同构成了企业核心竞争力的基础。"},
        {"heading": "🚀 行动路径", "content": "基于行业最佳实践，我们建议采取分阶段推进策略：第一阶段聚焦核心能力建设，第二阶段推动规模化增长，第三阶段实现生态协同。每个阶段有明确的目标、里程碑和成功标准，确保执行有方向、有节奏、有保障。"}
    ],
    "insight": "未来已来，关键是你如何应对这场变革",
    "takeaway": "明确方向 → 构建能力 → 规模发展 → 生态协同"
}
""".strip()

    try:
        api_format = MODEL_API_FORMATS.get(active_model, MODEL_API_FORMATS["openai"])
        model_name = os.environ.get(f"{active_model.upper()}_MODEL", model_config["default_model"])

        if api_format["format"] == "openai":
            headers = {
                "Authorization": f"{api_format['auth_header']} {api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "model": model_name,
                "messages": [
                    {"role": "system", "content": system},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7
            }
            resp = requests.post(model_config["api_url"], headers=headers, json=payload, timeout=60)

        elif api_format["format"] == "claude":
            headers = {
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "Content-Type": "application/json"
            }
            payload = {
                "model": model_name,
                "max_tokens": 2048,
                "messages": [{"role": "user", "content": prompt}]
            }
            resp = requests.post(model_config["api_url"], headers=headers, json=payload, timeout=60)

        elif api_format["format"] == "gemini":
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            payload = {
                "contents": [{"parts": [{"text": f"{system}\n\n{prompt}"}]}]
            }
            resp = requests.post(
                f"{model_config['api_url']}/{model_name}:generateContent",
                headers=headers, json=payload, timeout=60
            )
        elif api_format["format"] == "spark":
            headers = {"Content-Type": "application/json"}
            payload = {
                "header": {"app_id": api_key.split(":")[0] if ":" in api_key else api_key},
                "parameter": {"chat": {"domain": model_name, "temperature": 0.5}},
                "payload": {"message": {"text": [{"role": "user", "content": prompt}]}}
            }
            resp = requests.post(model_config["api_url"], headers=headers, json=payload, timeout=60)

        if resp.status_code == 200:
            if api_format["format"] == "gemini":
                return resp.json().get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
            elif api_format["format"] == "claude":
                return resp.json().get("content", [{}])[0].get("text", "")
            else:
                return resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        else:
            print(f"API调用失败: {resp.status_code} - {resp.text[:200]}")
            # FIX-P02: 非200状态码也视为失败，尝试切换到DeepSeek
            resp = None  # 标记需要切换

    except requests.exceptions.RequestException as e:
        print(f"AI API调用失败: {e}")
        resp = None  # 标记需要切换
    except Exception as e:
        print(f"AI API调用失败: {e}")
        resp = None  # 标记需要切换

    # FIX-P02: API调用失败时自动切换DeepSeek重试（仅当有有效的api_key时）
    if (resp is None or active_model == "minimax") and api_key and api_key != "fake-minimax-key":
        print(f"[FIX-P02] API调用失败或MiniMax模型，自动切换到DeepSeek...")
        try:
            fallback_model = "deepseek"
            fallback_config = AI_MODELS.get(fallback_model, AI_MODELS[DEFAULT_MODEL])
            fallback_env_key = fallback_config["env_key"]
            fallback_api_key = os.environ.get(fallback_env_key, "")

            if not fallback_api_key:
                fallback_api_key = os.environ.get("DEEPSEEK_API_KEY", "")

            if fallback_api_key:
                fallback_api_format = MODEL_API_FORMATS.get(fallback_model, MODEL_API_FORMATS["openai"])
                fallback_model_name = os.environ.get(f"{fallback_model.upper()}_MODEL", fallback_config["default_model"])

                headers = {
                    "Authorization": f"{fallback_api_format['auth_header']} {fallback_api_key}",
                    "Content-Type": "application/json"
                }
                payload = {
                    "model": fallback_model_name,
                    "messages": [
                        {"role": "system", "content": system},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.7
                }
                resp_fb = requests.post(fallback_config["api_url"], headers=headers, json=payload, timeout=60)

                if resp_fb.status_code == 200:
                    result = resp_fb.json()
                    print(f"[FIX-P02] DeepSeek切换成功")
                    return result.get("choices", [{}])[0].get("message", {}).get("content", "")
                else:
                    print(f"[FIX-P02] DeepSeek切换也失败: {resp_fb.status_code}")

        except Exception as e:
            print(f"[FIX-P02] DeepSeek切换异常: {e}")

    return ""

# ============================================================
# 行业专业Prompt系统
# ============================================================
INDUSTRY_PROMPTS = {
    "general": {
        "system": """你是顶级战略咨询顾问，擅长为CEO和高管生成决策级演示内容。
你的内容必须达到麦肯锡、贝恩、 BCG 的专业水准：
- 每个观点必须有具体数据支撑（百分比、金额、人数、时间）
- 每个结论必须有真实案例佐证（具体公司名称、行业、时间）
- 每个建议必须有可量化的预期效果

禁止：空洞的战略词汇、无数据支撑的判断、无法落地的方向性描述""",
        "content_template": """为PPT章节「{title}」生成高密度决策内容。

【Gamma级质量标准】
- 每个point必须包含：具体数字 + 真实案例 + 可操作建议
- 避免："提升效率"→改为"A公司通过X系统将某指标从Y%提升至Z%，周期T个月，投资回报R个月"
- 避免："市场增长"→改为"2024年市场规模N亿元，年增速X%，预计2026年达到Y亿元，CAGR Z%"

输出必须严格遵守此JSON格式（不要有任何额外内容）：
{{
    "title": "{title}",
    "subtitle": "精准定位的一句话描述",
    "points": [
        {{
            "heading": "【核心观点1】必须是一个完整的业务判断句",
            "content": "详细阐述：①具体数据（数字+单位+来源）②真实案例（公司名+做法+结果）③可操作建议（具体步骤+量化预期）。不少于80字。"
        }},
        {{
            "heading": "【核心观点2】必须是一个完整的业务判断句",
            "content": "详细阐述：①具体数据（数字+单位+来源）②真实案例（公司名+做法+结果）③可操作建议（具体步骤+量化预期）。不少于80字。"
        }},
        {{
            "heading": "【核心观点3】必须是一个完整的业务判断句",
            "content": "详细阐述：①具体数据（数字+单位+来源）②真实案例（公司名+做法+结果）③可操作建议（具体步骤+量化预期）。不少于80字。"
        }}
    ],
    "insight": "一句话核心洞察（必须包含具体数字或案例，不能是空洞金句）",
    "takeaway": "本章3个要点回顾（每个点10字内，用分号分隔）"
}}"""
    },
    "education": {
        "system": """你是教育行业顶级培训师，课程设计对标得到/混沌大学专业水准。
每个教学单元必须包含：
- 理论框架（认知科学/学习理论依据）
- 落地方法（具体操作步骤+工具推荐）
- 真实案例（学员数据+效果量化）
- 转化追踪（行为改变的可衡量指标）

禁止：空洞的教育理念、无数据支撑的教学承诺、无法评估的学习目标""",
        "content_template": """为教育培训PPT章节「{title}」生成专业教学设计内容。

【教学质量标准】
- 理论依据：引用具体研究（研究者+年份+核心结论）
- 方法设计：包含具体工具/步骤/时间分配
- 案例数据：必须是真实教学案例，包含学员数/满意度/行为改变率
- 效果追踪：必须有可量化的学习效果指标

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "学习目标导向的精准定位",
    "points": [
        {{
            "heading": "📚 理论框架（What）",
            "content": "核心概念定义+理论来源（研究者/年份/研究结论）+适用场景+使用限制。必须包含具体理论名称和学者名字。不少于60字。"
        }},
        {{
            "heading": "💡 实操方法（How）",
            "content": "具体操作步骤（编号列表）+工具/模板推荐+常见误区（3个）+技巧要点。必须可落地执行。不少于60字。"
        }},
        {{
            "heading": "🎯 案例应用（Example）",
            "content": "真实培训案例：企业/机构名+学员数量+教学设计+效果数据（满意度/行为改变率/绩效提升）。必须包含具体数字。不少于60字。"
        }}
    ],
    "quote": "一句教育名言+作者+应用场景（用于强化记忆）",
    "takeaway": "本章学习目标：完成此章节后学员能【具体做什么】，达到【什么可衡量标准】"
}}"""
    },
    "medical": {
        "system": """你是顶级医药战略顾问，内容对标IQVIA/艾昆纬研究报告水准。
医疗内容必须：
- 数据来源权威（NMPA/FDA/柳叶刀/JAMA等）
- 流行病学数据精确（发病率/患病率/死亡率/5年生存率）
- 治疗效果可量化（RR/OR/HR/P值/置信区间）
- 风险收益比清晰（必须包含NNT/NNH等指标）

禁止：未经证实的疗效声称、夸大治疗效果、回避严重不良反应""",
        "content_template": """为医疗健康PPT章节「{title}」生成专业临床导向内容。

【医学内容标准】
- 发病率/患病率：必须标注地区+年份+数据来源
- 治疗效果：必须包含相对风险降低（RRR）或绝对风险降低（ARR）+ NNT
- 药物数据：剂量+疗程+价格+医保情况
- 真实病例：脱敏后但保留关键临床特征

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "临床定位与流行病学背景",
    "points": [
        {{
            "heading": "📋 疾病认知与流行病学",
            "content": "定义+病因+诊断标准+关键流行病学数据（发病率/患病率/死亡率+地区+年份+数据来源）。必须包含具体数字。不少于60字。"
        }},
        {{
            "heading": "💊 治疗方案与药物经济学",
            "content": "一线/二线用药方案（通用名+商品名+剂量+疗程）+年治疗费用+医保目录情况+疗效数据（RRR/ARR/NNT）。不少于60字。"
        }},
        {{
            "heading": "🏥 临床实践与指南",
            "content": "最新临床指南推荐+真实世界研究数据+典型病例（关键体征+检验结果+治疗方案+转归）。必须脱敏但保留临床特征。不少于60字。"
        }}
    ],
    "stat": "关键流行病学数据（发病率/患病率/5年生存率，附来源）",
    "note": "参考文献或指南来源（格式：期刊名. 年份;卷(期):页码）",
    "takeaway": "本章临床要点：诊断要点+核心治疗原则+随访指标"
}}"""
    },
    "finance": {
        "system": """你是顶级投资银行分析师，内容对标高盛/摩根士丹利研报水准。
金融内容必须：
- 市场数据：包含具体金额单位（万亿/亿/百万）+ 数据来源
- 投资逻辑：明确的估值方法（PE/PB/PS/DCF）+ 具体假设
- 风险提示：量化风险敞口（VaR/最大回撤）+ 历史极端情景
- 配置建议：明确的仓位区间+再平衡触发条件

禁止：模糊的"高增长/低风险"描述、无具体价位的目标价、无时限的"长期看好" """,
        "content_template": """为金融投资PPT章节「{title}」生成专业投研内容。

【投研内容标准】
- 市场数据：规模（万亿/亿）+ 增速（%）+ 市场份额（头部三家占比）
- 投资分析：估值方法 + 核心假设 + 目标价位/区间 + 催化剂
- 风险收益：风险因素（已量化）+ 历史胜率/赔率 + 最大回撤历史
- 配置指引：明确仓位区间 + 标的筛选标准 + 买入/卖出触发条件

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "投资视角的核心定位",
    "points": [
        {{
            "heading": "📊 市场数据与规模",
            "content": "具体市场规模（万亿/亿+货币单位）+ 同比增速（%）+ 环比增速（%）+ 市场结构（CR3/CR5份额）+ 数据来源（Wind/Bloomberg/公司财报）。不少于80字。"
        }},
        {{
            "heading": "🎯 投资逻辑与估值",
            "content": "核心投资观点+估值方法（PE/PB/PS/DCF等）+ 核心假设+ 目标价位+ 时间轴+ 关键催化剂（政策/业绩/事件）。不少于80字。"
        }},
        {{
            "heading": "⚠️ 风险提示与应对",
            "content": "主要风险因素（已量化：最大回撤历史X%/ VaR X万）+ 历史情景（2008/2015/2020如何表现）+ 风险对冲策略+仓位管理建议。不少于80字。"
        }}
    ],
    "chart": {{"type": "bar", "title": "近三年核心财务指标对比", "categories": ["2022年", "2023年", "2024年"], "values": [具体数字1, 具体数字2, 具体数字3], "source": "公司财报/研究机构"}},
    "insight": "投资主线一句话：核心逻辑+目标涨幅+时间窗口",
    "takeaway": "配置建议：仓位区间+标的类型+再平衡条件（3句话以内）"
}}"""
    },
    "technology": {
        "system": """你是顶级科技行业分析师，内容对标Gartner/Forrester/IDC报告水准。
科技内容必须：
- 技术参数：具体版本号/性能指标/架构差异/兼容性
- 市场份额：头部三家份额（%）+ 增速对比 + 集中度变化趋势
- 竞品对比：具体功能差异表 + TCO对比 + 客户选择关键因素权重
- 技术路线：具体时间节点 + 技术成熟度 + 替代威胁

禁止：模糊的"AI赋能"描述、无具体参数的"性能提升"、无时间节点的"未来趋势" """,
        "content_template": """为科技互联网PPT章节「{title}」生成专业技术与市场分析内容。

【科技内容标准】
- 技术解读：具体架构/协议/版本号 + 性能基准测试数据
- 竞品分析：功能对比表 + TCO对比 + 市场份额具体数字
- 市场格局：CR3/CR5占比（%）+ 增速对比 + 集中度变化
- 趋势判断：具体时间节点（2025Q3/2026等）+ 技术成熟度评级

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "技术定位与市场格局",
    "points": [
        {{
            "heading": "🔬 技术原理与架构",
            "content": "核心技术原理+具体架构/协议名称+性能基准数据（具体数字+测试条件）+ 技术优势与局限性。必须包含具体参数。不少于80字。"
        }},
        {{
            "heading": "🏆 竞品对比与市场格局",
            "content": "头部三家市场份额（%名称+份额+增速）+ 核心差异对比表 + 客户选择因素权重（功能X%/价格Y%/服务Z%）+ 集中度变化趋势。不少于80字。"
        }},
        {{
            "heading": "🚀 技术路线与时间节点",
            "content": "未来18个月技术路线图（具体时间节点+里程碑）+ 技术成熟度预测（概念/实验/成熟/淘汰）+ 替代技术威胁评估 + 对现有产品影响。不少于80字。"
        }}
    ],
    "chart": {{"type": "bar", "title": "市场份额与增速对比", "categories": ["竞品A", "竞品B", "竞品C", "我们"], "values": [具体数字1, 具体数字2, 具体数字3, 具体数字4], "source": "IDC/Gartner 2024"}},
    "insight": "技术判断一句话：技术代际+窗口期+关键决定因素",
    "takeaway": "行动建议：技术选型建议+迁移时间窗口+风险预案（3句话）"
}}"""
    },
    "ecommerce": {
        "system": """你是顶级电商运营专家，内容对标天猫/京东/拼多多官方运营手册水准。
电商内容必须：
- 运营指标：具体数字+行业均值对比+优秀值标杆
- 策略拆解：具体步骤+工具选择+时间节点
- 案例数据：店铺名称+具体运营动作+GMV/转化率/复购率变化
- ROI测算：投入产出比+回本周期+盈亏平衡点

禁止：模糊的"提升转化"描述、无具体数字的"爆款打造"、无时间节点的"持续增长" """,
        "content_template": """为电商运营PPT章节「{title}」生成专业运营内容。

【电商内容标准】
- 核心指标：具体数字+行业均值+优秀值（TOP10%商家水平）
- 操作方法：具体步骤+工具推荐+操作时间节点
- 成功案例：店铺类型+具体运营动作+指标变化（前后对比）
- 投入产出：具体费用+预期产出+ROI+回本周期

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "运营视角与核心指标定位",
    "points": [
        {{
            "heading": "📈 核心指标与行业对标",
            "content": "关键运营指标（转化率/客单价/复购率/GMV+具体数字）+ 行业均值对比（行业均值X%/优秀值Y%）+ 优化空间分析。必须包含具体数字。不少于80字。"
        }},
        {{
            "heading": "🛠️ 执行方法与工具",
            "content": "具体操作步骤（编号列表）+ 工具/平台推荐 + 时间节点（什么时候做什么）+ 常见误区（3个）+ 操作注意事项。不少于80字。"
        }},
        {{
            "heading": "✅ 成功案例与效果",
            "content": "真实店铺案例：店铺类型/类目 + 具体运营动作 + 前后指标对比（GMV变化X%/转化率提升Y%/客单价提高Z元）+ 投入费用+ROI。不少于80字。"
        }}
    ],
    "chart": {{"type": "line", "title": "季度GMV与增速趋势", "categories": ["Q1", "Q2", "Q3", "Q4"], "values": [具体数字], "source": "平台数据/商家后台"}},
    "insight": "运营核心洞察一句话：关键成功因素+近期机会点",
    "takeaway": "可落地行动清单：本周做X件事 + 本月目标Y指标（3句话）"
}}"""
    },
    "government": {
        "system": """你是政府政策研究专家，内容对标国务院发展研究中心/国家发改委政策研究水准。
政务内容必须：
- 政策引用：具体文号+发布机关+发布时间+核心条款
- 数据来源：官方统计口径+年份+发布部门
- 实施路径：具体时间节点+责任部门+考核指标
- 预期成效：可量化指标+横向对比+历史经验

禁止：模糊的政策表述、无时间节点的"持续推进"、无责任主体的"加强领导" """,
        "content_template": """为政府企业PPT章节「{title}」生成规范性政策研究内容。

【政务内容标准】
- 政策依据：文件名称+文号+发布机关+发布时间+核心条款摘要
- 实施现状：具体数据+责任部门+时间节点+考核指标完成率
- 问题分析：具体表现+数据支撑+原因剖析
- 下步计划：具体任务+时间节点+预期成效+责任部门

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "政策背景与工作定位",
    "points": [
        {{
            "heading": "📜 政策背景与依据",
            "content": "相关政策文件名称+文号+发布机关+发布时间+核心条款摘要（具体条款号+内容）。政策背景（出台动机+现实需求）。不少于80字。"
        }},
        {{
            "heading": "📊 实施现状与数据",
            "content": "总体实施情况+具体数据（覆盖率/完成率/投入资金+单位）+ 责任部门+时间节点+考核指标+存在的主要问题（具体表现）。不少于80字。"
        }},
        {{
            "heading": "🎯 下步计划与成效",
            "content": "下步重点任务（具体工作+责任部门+时间节点）+ 预期目标（可量化指标）+ 保障措施+ 成功案例（哪个地区/部门+具体做法+成效数据）。不少于80字。"
        }}
    ],
    "insight": "政策精神一句话概括：核心目标+关键抓手+工作方向",
    "takeaway": "落实要点：本周/本月/本季度具体抓什么（3句话）"
}}"""
    },
    "realestate": {
        "system": """你是顶级房地产分析师，内容对标世联行/中指院/克而瑞研究报告水准。
房产内容必须：
- 项目分析：具体区位数据+配套成熟度+客群画像+竞品对比
- 市场数据：去化周期+存量面积+成交价格+环比同比变化
- 营销策略：具体定价策略+蓄客数据+开盘去化率+优惠力度
- 投资测算：IRR/NVP/回本周期+敏感性分析

禁止：模糊的"区位优势明显"、无具体价格的"性价比高"、无数据支撑的"热销" """,
        "content_template": """为房产建筑PPT章节「{title}」生成专业地产研究内容。

【地产内容标准】
- 项目禀赋：具体区位数据（交通/教育/医疗配套+距离/时间）+ 产品定位（面积段/总价段/客群画像）
- 市场数据：区域供销存（万平方米）+ 去化周期（月）+ 成交均价（元/㎡）+ 环比同比变化（%）
- 竞品分析：直接竞品项目+价格/产品/去化对比 + 差异化竞争优势
- 营销数据：蓄客量+认筹量+开盘去化率+来访成交比

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "项目定位与市场机会",
    "points": [
        {{
            "heading": "🏗️ 项目分析与产品定位",
            "content": "区位分析（具体位置+交通配套+教育医疗+距离）+ 产品定位（面积段+总价段+客群画像）+ 核心卖点（3个）+ 不利因素（坦诚告知）。不少于80字。"
        }},
        {{
            "heading": "📊 市场洞察与竞品",
            "content": "区域供销存数据（批准预售/成交/库存+万平方米+去化周期月）+ 成交均价（元/㎡+环比±%）+ 直接竞品对比（项目名+价格+产品+去化率）+ 差异化机会。不少于80字。"
        }},
        {{
            "heading": "💰 营销策略与定价",
            "content": "定价策略（均价+总价区间+与竞品差价）+ 营销节点（认筹/开盘时间）+ 优惠方案（具体折扣力度）+ 预期去化率+ 回本测算（IRR/NVP/回本周期）。不少于80字。"
        }}
    ],
    "chart": {{"type": "bar", "title": "区域价格对比（元/㎡）", "categories": ["本案", "竞品A", "竞品B", "区域均价"], "values": [具体数字], "source": "克而瑞 2024年Q3"}},
    "insight": "市场判断一句话：窗口期+风险点+核心机会",
    "takeaway": "操盘建议：定价区间+蓄客目标+开盘去化率目标（3句话）"
}}"""
    },
    "media": {
        "system": """你是顶级传媒广告策略顾问，内容对标奥美/阳狮/蓝色光标策略方案水准。
传媒内容必须：
- 受众洞察：具体画像数据+媒体习惯+内容偏好+消费场景
- 传播策略：具体媒介组合+内容形式+发布时间节奏
- 效果数据：曝光/互动/转化具体数字+行业均值对比+ROI
- 创意概念：核心创意阐述+执行要点+成功要素

禁止：模糊的"精准触达"、无具体数字的"刷屏"、无时间节点的"持续发酵" """,
        "content_template": """为传媒广告PPT章节「{title}」生成专业传播策略内容。

【传媒内容标准】
- 受众洞察：人口统计特征+媒体习惯（具体平台+使用时长）+ 内容偏好 + 消费决策路径
- 传播策略：媒介组合（具体平台+预算占比）+ 内容形式（具体）+ 发布时间节奏 + KOL/KOC策略
- 效果指标：预期曝光量/互动率/转化率（具体数字）+ 行业均值对比 + 成功案例参照
- 创意概念：核心创意主张+执行要点+差异化记忆点

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "内容主题与受众定位",
    "points": [
        {{
            "heading": "💡 受众洞察与画像",
            "content": "目标受众人口统计（年龄/职业/收入/城市）+ 媒体习惯（具体平台+使用时长）+ 内容偏好（类型/风格/话题）+ 消费决策路径（触点排序）。数据来源要标注。不少于80字。"
        }},
        {{
            "heading": "📢 传播策略与媒介组合",
            "content": "媒介组合策略（平台+预算占比+理由）+ 内容形式（长视频/短剧/图文/KOL具体类型）+ 发布时间节奏（预热/爆发/长尾）+ KOL/KOC具体名单或筛选标准。不少于80字。"
        }},
        {{
            "heading": "📈 效果预估与成功案例",
            "content": "预期效果指标（曝光X万/互动率X%/转化X单）+ 行业均值对比 + 具体成功案例（品牌+活动类型+效果数据）+ 风险预案（可能出现的问题+应对）。不少于80字。"
        }}
    ],
    "insight": "传播核心洞察一句话：关键成功因素+近期红利+差异化机会",
    "takeaway": "创意执行要点：核心记忆点+必做元素+规避事项（3句话）"
}}"""
    },
    "law": {
        "system": """你是顶级法律顾问，内容对标金杜/中伦/君合法律意见书水准。
法律内容必须：
- 法条引用：具体条款号+法条全文摘要+司法解释文号
- 案例数据：案件名称+审理法院+案号+裁判要点+裁判结果
- 风险评估：具体风险点+发生概率+法律后果+防范措施
- 合规建议：具体操作指引+所需材料+审批流程+时间节点

禁止：模糊的"存在法律风险"、无具体条款的"违反法律规定"、无应对方案的"需谨慎处理" """,
        "content_template": """为法律PPT章节「{title}」生成专业法律分析内容。

【法律内容标准】
- 法律依据：具体法条编号+条款内容摘要+司法解释文号+适用条件
- 案例分析：真实案例（案号+审理法院+裁判要点+裁判结果）
- 风险评估：具体风险点+司法实践中的认定+发生概率+法律后果
- 合规建议：具体操作建议+所需材料+审批流程+时间节点+费用预估

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "法律场景与风险定位",
    "points": [
        {{
            "heading": "⚖️ 法律依据与适用",
            "content": "相关法律条款（名称+具体条款号+条款内容摘要）+ 司法解释（文号+核心内容）+ 适用条件+ 注意事项（实践中的争议焦点）。法条引用必须准确。不少于80字。"
        }},
        {{
            "heading": "📋 典型案例与裁判要点",
            "content": "真实案例分析：案件名称+审理法院+案号+当事人+核心争议焦点+法院裁判要点+裁判结果+ 实务启示。案例要有代表性。不少于80字。"
        }},
        {{
            "heading": "🛡️ 风险防范与合规建议",
            "content": "主要风险点（具体表现+发生概率+法律后果）+ 风险防范措施（具体操作步骤）+ 合规建议（所需材料+审批流程+时间节点+费用）。必须可落地执行。不少于80字。"
        }}
    ],
    "insight": "法律核心提示一句话：最关键风险点+最重要防范措施",
    "takeaway": "本章实务要点：做什么+怎么做+注意什么（3句话）"
}}"""
    },
    "manufacture": {
        "system": """你是顶级制造业咨询顾问，内容对标罗兰贝格/波士顿咨询制造业报告水准。
制造业内容必须：
- 产能数据：具体产能规模+产能利用率+扩产投资+工艺路线
- 效率指标：OEE/良品率/单位能耗/人均产值具体数字+行业对标
- 供应链：具体供应商名单+采购占比+替代风险+物流成本
- 技术升级：具体设备名称+投资额+回报周期+技术差距

禁止：模糊的"产能充足"、无具体数字的"效率提升"、无供应商名称的"供应链稳定" """,
        "content_template": """为制造业PPT章节「{title}」生成专业生产运营内容。

【制造业内容标准】
- 产能与工艺：具体产能规模（万件/年）+ 产能利用率（%）+ 核心工艺路线 + 扩产计划（投资+时间节点）
- 效率指标：OEE/良品率/单位能耗/人均产值（具体数字）+ 行业对标（行业均值/先进值）+ 差距分析
- 供应链：核心原材料供应商（名称+占比%）+ 替代供应商情况 + 物流成本占比 + 供应链风险评估
- 技术升级：具体设备/系统名称+投资额+预期效果+回本周期

输出JSON格式（必须严格遵守）：
{{
    "title": "{title}",
    "subtitle": "生产运营与效率定位",
    "points": [
        {{
            "heading": "🏭 产能布局与工艺水平",
            "content": "产能规模（具体数字+单位）+ 产能利用率（%）+ 核心工艺路线+设备水平 + 扩产计划（具体投资+时间节点+预期新增产能）。不少于80字。"
        }},
        {{
            "heading": "📊 效率指标与行业对标",
            "content": "OEE/良品率/单位能耗/人均产值（具体数字）+ 行业标杆对比（国际龙头/国内龙头具体数字）+ 差距根本原因分析 + 提升路径。不少于80字。"
        }},
        {{
            "heading": "🔧 技术升级与投资回报",
            "content": "待升级具体设备/系统+投资额+预期效果（良品率+效率+成本具体变化）+ 回本周期+ 实施风险+ 替代进口设备情况。不少于80字。"
        }}
    ],
    "chart": {{"type": "bar", "title": "效率指标行业对标", "categories": ["我们", "行业均值", "行业先进"], "values": [具体数字], "source": "行业协会 2024"}},
    "insight": "制造核心洞察一句话：最关键差距+最快见效举措+长期竞争力来源",
    "takeaway": "行动建议：优先升级什么+投入多少+预期回报（3句话）"
}}"""
    }
}


def get_industry_prompt(industry: str, chapter_title: str) -> tuple:
    """获取行业专属的system prompt和content template"""
    industry_key = industry if industry in INDUSTRY_PROMPTS else "general"
    config = INDUSTRY_PROMPTS[industry_key]
    template = config["content_template"].format(title=chapter_title)
    return config["system"], template

def parse_ai_content(ai_output: str) -> dict:
    """解析AI输出的JSON内容，提取结构化数据"""
    try:
        match = re.search(r'\{.*\}', ai_output, re.DOTALL)
        if match:
            data = json.loads(match.group())
            # 确保chart数据格式正确
            if "chart" in data and isinstance(data["chart"], dict):
                # 确保values是数字列表
                if "values" in data["chart"]:
                    data["chart"]["values"] = [
                        int(v) if str(v).isdigit() else float(v)
                        for v in data["chart"]["values"]
                    ]
            return data
    except Exception as e:
        print(f"JSON解析失败: {e}")

    # 降级处理：尝试从纯文本中提取结构
    lines = ai_output.strip().split('\n')
    points = []
    current_heading = ""
    current_content = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        # 检测标题行（通常是加粗、序号、或者短句）
        if len(line) < 50 and (line.startswith(('1.', '2.', '3.', '•', '-', '▸', '●')) or line.isupper()):
            if current_heading:
                points.append({"heading": current_heading, "content": " ".join(current_content)})
            current_heading = line.lstrip('123456789.-•▸● ')
            current_content = []
        elif current_heading:
            current_content.append(line)

    if current_heading:
        points.append({"heading": current_heading, "content": " ".join(current_content)})

    if not points:
        # 终极降级：整段作为内容
        return {
            "title": "",
            "points": [{"heading": "核心内容", "content": ai_output[:500]}]
        }

    return {"points": points}

# ============================================================
# 专业模板渲染系统 - Gamma级专业设计
# ============================================================
TEMPLATE_LAYOUTS = {
    # ========== Gamma风格通用配置 ==========
    # 核心设计原则：大标题+卡片布局+充足留白+专业色彩
    "academic": {
        "name": "学术风",
        "title_font": 44,
        "content_font": 18,
        "heading_font": 22,
        "left_bar_width": 0.6,
        "card_bg": (248, 250, 252),
        "card_border": (226, 232, 240),
        "accent_color": (44, 82, 130),
        "title_color": (26, 54, 93),
        "decorative": "line_bottom",
        "layout": "cards_3col",
    },
    "business": {
        "name": "商务蓝",
        "title_font": 48,
        "content_font": 20,
        "heading_font": 26,
        "left_bar_width": 0.7,
        "card_bg": (248, 250, 252),
        "card_border": (203, 213, 225),
        "accent_color": (237, 137, 54),
        "title_color": (26, 54, 93),
        "decorative": "corner_block",
        "layout": "cards_2col",
    },
    "enterprise": {
        "name": "企业蓝",
        "title_font": 44,
        "content_font": 18,
        "heading_font": 24,
        "left_bar_width": 0.5,
        "card_bg": (241, 245, 249),
        "card_border": (203, 213, 225),
        "accent_color": (49, 130, 206),
        "title_color": (26, 54, 93),
        "decorative": "left_bar",
        "layout": "cards_2col",
    },
    "tech": {
        "name": "科技风",
        "title_font": 52,
        "content_font": 20,
        "heading_font": 28,
        "left_bar_width": 0.8,
        "card_bg": (15, 23, 42),
        "card_border": (51, 65, 85),
        "accent_color": (0, 212, 255),
        "title_color": (15, 23, 42),
        "decorative": "gradient_side",
        "layout": "cards_2col",
        "dark_mode": True,
    },
    "simple": {
        "name": "简约风",
        "title_font": 40,
        "content_font": 16,
        "heading_font": 20,
        "left_bar_width": 0,
        "card_bg": (255, 255, 255),
        "card_border": (241, 245, 249),
        "accent_color": (100, 116, 139),
        "title_color": (30, 41, 59),
        "decorative": "none",
        "layout": "bullets",
    },
    "gradient": {
        "name": "渐变风",
        "title_font": 50,
        "content_font": 20,
        "heading_font": 26,
        "left_bar_width": 0,
        "card_bg": (255, 255, 255),
        "card_border": (226, 232, 240),
        "accent_color": (99, 102, 241),
        "title_color": (30, 41, 59),
        "decorative": "gradient_bg",
        "layout": "big_number_cards",
    },
    "default": {
        "name": "默认",
        "title_font": 44,
        "content_font": 18,
        "heading_font": 22,
        "left_bar_width": 0.5,
        "card_bg": (248, 250, 252),
        "card_border": (226, 232, 240),
        "accent_color": (44, 82, 130),
        "title_color": (26, 54, 93),
        "decorative": "none",
        "layout": "cards_2col",
    }
}

def get_template_config(template_id: str) -> dict:
    """获取模板配置"""
    return TEMPLATE_LAYOUTS.get(template_id, TEMPLATE_LAYOUTS["default"])

def apply_template_compatibility(config: dict) -> dict:
    """将新格式模板配置适配到旧函数兼容格式"""
    return {
        "title_font": config.get("title_font", 44),
        "content_font": config.get("content_font", 18),
        "heading_font": config.get("heading_font", 22),
        "accent_width": config.get("left_bar_width", 0.5),
        "has_header_bar": config.get("decorative", "none") != "none",
        "header_height": 0.8 if config.get("has_header_bar") else 0,
        "decorative_shapes": config.get("decorative", "none"),
        "bullet_style": "filled_circle",
    }

def add_decorative_shapes(slide, shapes, template_id: str, colors: dict, page_type: str = "content"):
    """为幻灯片添加装饰元素"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = shapes.parent  # 获取presentation对象

    if template_id in ["tech", "cyber", "future"]:
        # 科技风格：添加角落装饰块
        # 左上角装饰
        left_accent = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), Inches(1.5)
        )
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = RGBColor(*colors.get("accent", (0, 212, 255)))
        left_accent.line.fill.background()

        # 底部装饰线
        bottom_line = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.2),
            Inches(13.333), Inches(0.05)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = RGBColor(*colors.get("accent", (0, 212, 255)))
        bottom_line.line.fill.background()

    elif template_id == "business":
        # 商务风格：右下角装饰
        corner = shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(11.5), Inches(5.5),
            Inches(1.8), Inches(2)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = RGBColor(*colors.get("accent", (237, 137, 54)))
        corner.line.fill.background()
        corner.rotation = 180

    elif template_id == "enterprise":
        # 企业风格：左侧竖条
        left_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = RGBColor(*colors.get("accent", (49, 130, 206)))
        left_bar.line.fill.background()

    # 通用：标题下划线装饰
    if page_type != "cover":
        try:
            line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(1.2),
                Inches(2), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
            line.line.fill.background()
        except:
            pass

def add_header_bar(slide, shapes, colors: dict, config: dict):
    """添加页面顶部装饰条"""
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    if config.get("has_header_bar"):
        bar_height = config.get("header_height", 0.8)
        header = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(bar_height)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        header.line.fill.background()

def add_image_to_slide(slide, image_path: str, left: float, top: float, width: float, height: float):
    """为幻灯片添加图片（带圆角遮罩）"""
    from pptx.util import Inches
    from PIL import Image
    import os
    
    if not image_path or not os.path.exists(image_path):
        return
    
    try:
        # 添加图片
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
        
        # 给图片加一层淡淡的白色遮罩（让图片与PPT风格融合）
        # 注意：python-pptx不支持直接给图片加滤镜，用半透明形状模拟
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
        overlay.fill.fore_color.brightness = 0.1  # 10%透明度效果（python-pptx不支持真透明，用亮度模拟）
        overlay.line.fill.background()
        # 压到图片下面（调整z-order，但python-pptx的z-order就是添加顺序）
    except Exception as e:
        pass

def add_gamma_chart(slide, chart_data: dict, colors: dict, left: float, top: float, width: float, height: float):
    """用形状绘制Gamma级精美图表（柱状/折线/饼图）"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    chart_type = chart_data.get("type", "bar")
    categories = chart_data.get("categories", [])
    values = chart_data.get("values", [])
    title = chart_data.get("title", "")

    if not categories or not values:
        return

    # Gamma配色方案
    palette = [
        colors.get("accent", (44, 82, 130)),
        colors.get("title", (26, 54, 93)),
        (99, 102, 241),   # 紫色
        (0, 212, 255),     # 青色
        (52, 211, 153),    # 绿色
        (251, 146, 60),    # 橙色
    ]

    # 标题
    if title:
        title_box = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        p.alignment = PP_ALIGN.LEFT

    chart_top = top + 0.1
    chart_height = height - 0.4
    bar_width = width / len(categories) * 0.6
    bar_gap = width / len(categories) * 0.4

    if chart_type == "bar":
        max_val = max(values) if values else 1
        for i, (cat, val) in enumerate(zip(categories, values)):
            bar_h = (val / max_val) * chart_height * 0.75
            x = left + i * (bar_width + bar_gap) + bar_gap / 2
            y = chart_top + chart_height - bar_h

            # 圆角矩形柱
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(bar_width), Inches(bar_h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*palette[i % len(palette)])
            bar.line.fill.background()

            # 数值标签
            val_box = slide.shapes.add_textbox(Inches(x), Inches(y - 0.25), Inches(bar_width), Inches(0.25))
            vt = val_box.text_frame
            vt.text = str(val)
            vp = vt.paragraphs[0]
            vp.font.size = Pt(10)
            vp.font.bold = True
            vp.font.color.rgb = RGBColor(*palette[i % len(palette)])
            vp.alignment = PP_ALIGN.CENTER

            # 分类标签
            cat_box = slide.shapes.add_textbox(Inches(x), Inches(chart_top + chart_height + 0.05), Inches(bar_width), Inches(0.3))
            ct = cat_box.text_frame
            ct.text = str(cat)
            cp = ct.paragraphs[0]
            cp.font.size = Pt(9)
            cp.font.color.rgb = RGBColor(100, 100, 100)
            cp.alignment = PP_ALIGN.CENTER

    elif chart_type == "line":
        max_val = max(values) if values else 1
        min_val = min(values) if values else 0
        range_val = max_val - min_val if max_val != min_val else 1
        step_x = width / (len(values) - 1) if len(values) > 1 else width

        for i, val in enumerate(values):
            x = left + i * step_x
            y = chart_top + chart_height - ((val - min_val) / range_val) * chart_height * 0.75

            # 数据点圆
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.08), Inches(y - 0.08),
                Inches(0.16), Inches(0.16)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*palette[0])
            dot.line.fill.background()

            # 数值标签
            val_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y - 0.3), Inches(0.6), Inches(0.25))
            vt = val_box.text_frame
            vt.text = str(val)
            vp = vt.paragraphs[0]
            vp.font.size = Pt(9)
            vp.font.bold = True
            vp.font.color.rgb = RGBColor(*palette[0])
            vp.alignment = PP_ALIGN.CENTER

            # 连线
            if i < len(values) - 1:
                next_x = left + (i + 1) * step_x
                next_y = chart_top + chart_height - ((values[i + 1] - min_val) / range_val) * chart_height * 0.75
                line = slide.shapes.add_connector(
                    1,  # STRAIGHT
                    Inches(x), Inches(y),
                    Inches(next_x), Inches(next_y)
                )
                line.line.color.rgb = RGBColor(*palette[0])
                line.line.width = Pt(2)

        # X轴标签
        for i, cat in enumerate(categories):
            x = left + i * step_x
            cat_box = slide.shapes.add_textbox(Inches(x - 0.4), Inches(chart_top + chart_height + 0.05), Inches(0.8), Inches(0.3))
            ct = cat_box.text_frame
            ct.text = str(cat)
            cp = ct.paragraphs[0]
            cp.font.size = Pt(9)
            cp.font.color.rgb = RGBColor(100, 100, 100)
            cp.alignment = PP_ALIGN.CENTER

    elif chart_type == "pie":
        total = sum(values)
        if total == 0:
            return
        angles = [v / total * 360 for v in values]
        center_x = left + width / 2
        center_y = chart_top + chart_height / 2
        radius = min(width, chart_height) / 2 * 0.85

        # 绘制饼图 - 用弧形近似
        start_angle = 0
        for i, (angle, cat, val) in enumerate(zip(angles, categories, values)):
            # 简化：用椭圆和三角形绘制饼块
            pie = slide.shapes.add_shape(
                MSO_SHAPE.PIE,
                Inches(center_x - radius), Inches(center_y - radius),
                Inches(radius * 2), Inches(radius * 2)
            )
            pie.fill.solid()
            pie.fill.fore_color.rgb = RGBColor(*palette[i % len(palette)])
            pie.line.fill.background()

            # 百分比标签
            pct = int(val / total * 100) if total > 0 else 0
            label_angle = start_angle + angle / 2
            label_x = center_x + (radius + 0.3) * 0.7 * (1 if label_angle < 180 else -1)
            label_y = center_y - (radius + 0.3) * 0.5

            pct_box = slide.shapes.add_textbox(Inches(label_x - 0.4), Inches(label_y), Inches(0.8), Inches(0.3))
            pt = pct_box.text_frame
            pt.text = f"{pct}%"
            pp = pt.paragraphs[0]
            pp.font.size = Pt(10)
            pp.font.bold = True
            pp.font.color.rgb = RGBColor(*palette[i % len(palette)])
            pp.alignment = PP_ALIGN.CENTER

            start_angle += angle

        # 图例
        legend_y = chart_top + chart_height + 0.25
        for i, cat in enumerate(categories):
            lx = left + (i % 4) * (width / 4)
            ly = legend_y if i < 4 else legend_y + 0.25

            dot = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(lx), Inches(ly),
                Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*palette[i % len(palette)])
            dot.line.fill.background()

            leg_box = slide.shapes.add_textbox(Inches(lx + 0.2), Inches(ly - 0.02), Inches(width / 4 - 0.3), Inches(0.25))
            lt = leg_box.text_frame
            lt.text = str(cat)
            lp = lt.paragraphs[0]
            lp.font.size = Pt(8)
            lp.font.color.rgb = RGBColor(80, 80, 80)


def add_chart_to_slide(slide, chart_data: dict, colors: dict, left: float, top: float, width: float, height: float):
    """为幻灯片添加图表（使用Gamma级形状图表）"""
    # 使用Gamma级形状图表
    add_gamma_chart(slide, chart_data, colors, left, top, width, height)

def render_structured_slide(slide, page_data: dict, colors: dict, template_id: str, page_type: str = "content"):
    """渲染结构化内容幻灯片 - Gamma级专业视觉设计"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    config = get_template_config(template_id)
    # 兼容旧函数格式
    compat = apply_template_compatibility(config)
    shapes = slide.shapes

    title = page_data.get("title", "")
    subtitle = page_data.get("subtitle", "")
    points = page_data.get("points", [])
    chart = page_data.get("chart")
    insight = page_data.get("insight", "")
    takeaway = page_data.get("takeaway", "")

    # ========== 封面页 ==========
    if page_type == "cover":
        top_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(4.5)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        top_bar.line.fill.background()

        accent_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.25), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
        accent_bar.line.fill.background()

        title_box = shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if subtitle:
            sub_box = shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(11), Inches(0.8))
            stf = sub_box.text_frame
            stf.text = subtitle
            sp = stf.paragraphs[0]
            sp.font.size = Pt(22)
            sp.font.color.rgb = RGBColor(200, 220, 255)

        info_box = shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11), Inches(0.6))
        inf_tf = info_box.text_frame
        inf_tf.text = "Powered by DemoPPT"
        inf_p = inf_tf.paragraphs[0]
        inf_p.font.size = Pt(14)
        inf_p.font.color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))

        bottom_line = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(5.7),
            Inches(3), Inches(0.06)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
        bottom_line.line.fill.background()

        corner_block = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(11), Inches(5.5),
            Inches(2.333), Inches(2)
        )
        corner_block.fill.solid()
        corner_block.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
        corner_block.line.fill.background()
        return

    # ========== 目录页 ==========
    if page_type == "toc":
        left_block = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(4.5), Inches(7.5)
        )
        left_block.fill.solid()
        left_block.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        left_block.line.fill.background()

        toc_title = shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5))
        toc_tf = toc_title.text_frame
        toc_tf.word_wrap = True
        toc_tf.text = "CONTENTS"
        toc_p = toc_tf.paragraphs[0]
        toc_p.font.size = Pt(42)
        toc_p.font.bold = True
        toc_p.font.color.rgb = RGBColor(255, 255, 255)

        num_box = shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(3.5), Inches(0.5))
        num_tf = num_box.text_frame
        num_tf.text = f"{len(points)} 个章节"
        num_p = num_tf.paragraphs[0]
        num_p.font.size = Pt(16)
        num_p.font.color.rgb = RGBColor(180, 200, 240)

        start_y = 1.0
        for i, point in enumerate(points[:8]):
            item_title = point.get("heading", point.get("title", f"章节{i+1}"))

            num_circle = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(5.2), Inches(start_y + i * 0.8),
                Inches(0.35), Inches(0.35)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
            num_circle.line.fill.background()

            num_txt = shapes.add_textbox(Inches(5.2), Inches(start_y + i * 0.8 - 0.02), Inches(0.35), Inches(0.35))
            nt = num_txt.text_frame
            nt.paragraphs[0].alignment = PP_ALIGN.CENTER
            nt.text = str(i + 1)
            np = nt.paragraphs[0]
            np.font.size = Pt(11)
            np.font.bold = True
            np.font.color.rgb = RGBColor(255, 255, 255)

            item_box = shapes.add_textbox(Inches(5.7), Inches(start_y + i * 0.8), Inches(7), Inches(0.4))
            ib = item_box.text_frame
            ib.text = item_title
            ip = ib.paragraphs[0]
            ip.font.size = Pt(16)
            ip.font.color.rgb = RGBColor(51, 51, 51)

            if i < len(points) - 1:
                line = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(5.37), Inches(start_y + i * 0.8 + 0.35),
                    Inches(0.02), Inches(0.45)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(220, 220, 220)
                line.line.fill.background()
        return

    # ========== 总结页 ==========
    if page_type == "summary":
        top_bar = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        top_bar.line.fill.background()

        title_box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title or "总结"
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

        card_top = 1.5
        for i, point in enumerate(points[:4]):
            card_bg = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5 + (i % 2) * 6.3), Inches(card_top + (i // 2) * 2.2),
                Inches(6.0), Inches(2.0)
            )
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
            card_bg.line.color.rgb = RGBColor(226, 232, 240)

            num_lbl = shapes.add_textbox(Inches(0.7 + (i % 2) * 6.3), Inches(card_top + 0.15 + (i // 2) * 2.2), Inches(0.5), Inches(0.4))
            nl = num_lbl.text_frame
            nl.text = f"0{i+1}"
            nlp = nl.paragraphs[0]
            nlp.font.size = Pt(24)
            nlp.font.bold = True
            nlp.font.color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))

            card_title = shapes.add_textbox(Inches(1.3 + (i % 2) * 6.3), Inches(card_top + 0.2 + (i // 2) * 2.2), Inches(5.0), Inches(0.4))
            ct = card_title.text_frame
            ct.text = point.get("heading", "")
            ctp = ct.paragraphs[0]
            ctp.font.size = Pt(16)
            ctp.font.bold = True
            ctp.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

            card_content = shapes.add_textbox(Inches(0.7 + (i % 2) * 6.3), Inches(card_top + 0.7 + (i // 2) * 2.2), Inches(5.5), Inches(1.2))
            cc = card_content.text_frame
            cc.word_wrap = True
            cc.text = point.get("content", "")[:120]
            ccp = cc.paragraphs[0]
            ccp.font.size = Pt(12)
            ccp.font.color.rgb = RGBColor(100, 100, 100)

        thanks_box = shapes.add_textbox(Inches(0), Inches(6.5), Inches(13.333), Inches(0.8))
        th = thanks_box.text_frame
        th.paragraphs[0].alignment = PP_ALIGN.CENTER
        th.text = "感谢观看  |  Thank You"
        thp = th.paragraphs[0]
        thp.font.size = Pt(28)
        thp.font.bold = True
        thp.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        return

    # ========== 内容页（默认） ==========
    left_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.12), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
    left_bar.line.fill.background()

    top_line = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.03)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
    top_line.line.fill.background()

    title_box = shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(compat.get("title_font", 36))
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

    start_y = 1.2
    if subtitle:
        sub_box = shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.4))
        stf = sub_box.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
        start_y = 1.5

    if insight:
        insight_box = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(start_y),
            Inches(12), Inches(0.5)
        )
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))
        insight_box.line.fill.background()

        ins_text = shapes.add_textbox(Inches(0.7), Inches(start_y + 0.08), Inches(11.6), Inches(0.4))
        it = ins_text.text_frame
        it.text = f"INSIGHT: {insight}"
        ip = it.paragraphs[0]
        ip.font.size = Pt(13)
        ip.font.italic = True
        ip.font.color.rgb = RGBColor(255, 255, 255)
        start_y += 0.7

    if chart and len(points) >= 1:
        col_width = 5.5
        chart_left = 7.0
        chart_top = start_y
        chart_width = 5.8
        chart_height = 4.5

        for i, point in enumerate(points[:2]):
            pt = start_y + i * 2.2

            num_shape = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(pt),
                Inches(0.35), Inches(0.35)
            )
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
            num_shape.line.fill.background()

            num_t = shapes.add_textbox(Inches(0.5), Inches(pt - 0.02), Inches(0.35), Inches(0.35))
            nt = num_t.text_frame
            nt.paragraphs[0].alignment = PP_ALIGN.CENTER
            nt.text = str(i + 1)
            np = nt.paragraphs[0]
            np.font.size = Pt(12)
            np.font.bold = True
            np.font.color.rgb = RGBColor(255, 255, 255)

            h_box = shapes.add_textbox(Inches(1.0), Inches(pt - 0.05), Inches(col_width - 0.5), Inches(0.4))
            ht = h_box.text_frame
            ht.text = point.get("heading", "")
            hp = ht.paragraphs[0]
            hp.font.size = Pt(compat.get("heading_font", 18))
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

            c_box = shapes.add_textbox(Inches(1.0), Inches(pt + 0.4), Inches(col_width - 0.5), Inches(1.8))
            ct = c_box.text_frame
            ct.word_wrap = True
            ct.text = point.get("content", "")
            cp = ct.paragraphs[0]
            cp.font.size = Pt(compat.get("content_font", 14))
            cp.font.color.rgb = RGBColor(55, 65, 81)
            cp.line_spacing = 1.35

        add_chart_to_slide(slide, chart, colors, chart_left, chart_top, chart_width, chart_height)

    elif len(points) >= 2:
        cols = 2 if len(points) <= 4 else 3
        col_w = 5.8 if cols == 2 else 4.0
        row_h = 2.3

        for i, point in enumerate(points[:cols * 2]):
            col = i % cols
            row = i // cols
            left = 0.5 + col * (col_w + 0.3)
            top = start_y + row * (row_h + 0.3)

            card = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top),
                Inches(col_w), Inches(row_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(248, 250, 252)
            card.line.color.rgb = RGBColor(226, 232, 240)

            top_accent = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top),
                Inches(col_w), Inches(0.08)
            )
            top_accent.fill.solid()
            top_accent.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
            top_accent.line.fill.background()

            num_t = shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.2), Inches(0.5), Inches(0.4))
            nt = num_t.text_frame
            nt.text = f"0{i+1}"
            np = nt.paragraphs[0]
            np.font.size = Pt(20)
            np.font.bold = True
            np.font.color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))

            h_box = shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.65), Inches(col_w - 0.3), Inches(0.4))
            ht = h_box.text_frame
            ht.word_wrap = True
            ht.text = point.get("heading", "")
            hp = ht.paragraphs[0]
            hp.font.size = Pt(15)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

            c_box = shapes.add_textbox(Inches(left + 0.15), Inches(top + 1.1), Inches(col_w - 0.3), Inches(1.1))
            ct = c_box.text_frame
            ct.word_wrap = True
            ct.text = point.get("content", "")[:150]
            cp = ct.paragraphs[0]
            cp.font.size = Pt(11)
            cp.font.color.rgb = RGBColor(55, 65, 81)
            cp.line_spacing = 1.3

    else:
        for i, point in enumerate(points):
            pt = start_y + i * 1.6

            dot = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(pt + 0.08),
                Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
            dot.line.fill.background()

            h_box = shapes.add_textbox(Inches(0.75), Inches(pt), Inches(11.5), Inches(0.4))
            ht = h_box.text_frame
            ht.text = point.get("heading", "")
            hp = ht.paragraphs[0]
            hp.font.size = Pt(compat.get("heading_font", 18))
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(*colors.get("title", (26, 54, 93)))

            c_box = shapes.add_textbox(Inches(0.75), Inches(pt + 0.45), Inches(11.5), Inches(1.0))
            ct = c_box.text_frame
            ct.word_wrap = True
            ct.text = point.get("content", "")
            cp = ct.paragraphs[0]
            cp.font.size = Pt(compat.get("content_font", 14))
            cp.font.color.rgb = RGBColor(55, 65, 81)
            cp.line_spacing = 1.4

    if takeaway:
        tw_box = shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
        tw = tw_box.text_frame
        tw.text = f"TAKEAWAY: {takeaway}"
        twp = tw.paragraphs[0]
        twp.font.size = Pt(11)
        twp.font.color.rgb = RGBColor(*colors.get("accent", (44, 82, 130)))
    else:
        footer_box = shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        ftf = footer_box.text_frame
        ftf.text = "DemoPPT"
        fp = ftf.paragraphs[0]
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(180, 180, 180)




@app.get("/api/set_model")
async def set_model(model: str = None):
    """设置默认AI模型"""
    if model and model in AI_MODELS:
        os.environ["ACTIVE_MODEL"] = model
        return {"success": True, "message": f"已切换到 {AI_MODELS[model]['name']}"}
    return {"success": False, "message": "模型不存在"}

@app.get("/")
async def root():
    return {"message": "DemoPPT API服务运行中", "version": "2.3.0"}

@app.get("/health")
async def health():
    return {"status": "ok"}

@app.post("/api/generate_outline")
def generate_outline(req: OutlineRequest, authorization: str = Header(None)):
    """生成PPT大纲"""
    # 从Header获取用户信息
    from main import verify_token
    user = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization[7:]
        user = verify_token(token)
    user_id = user["user_id"] if user else 0

    industry = req.industry or "general"
    industry_info = INDUSTRY_TEMPLATES.get(industry, INDUSTRY_TEMPLATES["general"])

    # 1. 从用户知识库召回相关内容
    kb_context = kb.get_context_for_prompt(req.topic, industry, top_k=3)

    # 2. 从预置行业知识库获取行业数据/洞察
    industry_context_data = search_and_get_context(req.topic, industry)
    industry_context = industry_context_data.get("full_context", get_industry_context(req.topic, industry))
    has_web_data = industry_context_data.get("has_web", False)

    # 构建增强上下文
    context_parts = []
    if kb_context:
        context_parts.append(f"【用户知识库参考】：\n{kb_context}")
    context_parts.append(f"【行业知识库】：\n{industry_context}")

    full_context = "\n\n".join(context_parts)

    # 获取用户配置的模型（用真正的user_id）
    user_api_config = None
    if req.model_config_id:
        from model_config import get_config_by_id, get_default_config
        user_api_config = get_config_by_id(req.model_config_id, user_id=user_id)
        if not user_api_config:
            user_api_config = get_default_config(user_id=user_id)
    else:
        from model_config import get_default_config
        user_api_config = get_default_config(user_id=user_id)

    # 未登录且无有效配置时返回错误提示
    if not user and not user_api_config:
        return {"success": False, "error": "请先登录或配置AI模型"}

    system_prompt = f"""你是PPT策划专家，擅长为{industry_info['name']}领域生成逻辑清晰的大纲。
要求：结构严谨、层次分明、符合行业规范。

{full_context}"""

    prompt = f"""请为主题「{req.topic}」生成{industry_info['name']}领域的PPT大纲。

{full_context}

要求：
- 5-8页幻灯片
- 包含封面、目录、内容章节、总结
- 每个章节标题简洁有力

返回JSON格式：
{{"pages": [{{"title": "标题", "type": "cover/toc/content/summary"}}]}}"""

    result = call_ai_model(prompt, system=system_prompt, user_api_config=user_api_config)

    try:
        match = re.search(r'\{.*\}', result, re.DOTALL)
        if match:
            data = json.loads(match.group())
            outline = data.get("pages", [])
            if outline:
                return {"success": True, "outline": outline, "has_kb_context": bool(kb_context), "has_web_data": has_web_data}
    except:
        pass

    # 使用高质量降级大纲
    outline = get_fallback_outline(req.topic, industry)
    return {"success": True, "outline": outline, "has_kb_context": False}


# ===== 内容预览 =====
class PreviewRequest(BaseModel):
    topic: str
    subtitle: Optional[str] = ""
    industry: Optional[str] = "general"
    page_title: str
    page_type: str = "content"
    model_config_id: Optional[int] = None  # 用户配置的模型ID

@app.post("/api/preview_content")
async def preview_content(req: PreviewRequest):
    """单页内容预览（不生成完整PPT，只预览该页AI会写什么）"""
    industry = req.industry or "general"
    industry_info = INDUSTRY_TEMPLATES.get(industry, INDUSTRY_TEMPLATES["general"])

    kb_context = kb.get_context_for_prompt(req.topic, industry, top_k=2)
    industry_context_data = search_and_get_context(req.topic, industry)
    industry_context = industry_context_data.get("full_context", get_industry_context(req.topic, industry))

    context_parts = []
    if kb_context:
        context_parts.append(f"【用户知识库参考】：\n{kb_context}")
    context_parts.append(f"【行业知识库】：\n{industry_context}")
    full_context = "\n\n".join(context_parts)

    # 获取用户配置的模型
    user_api_config = None
    if req.model_config_id:
        from model_config import get_config_by_id, get_default_config
        user_api_config = get_config_by_id(req.model_config_id, user_id=0)
        if not user_api_config:
            user_api_config = get_default_config(user_id=0)
    else:
        from model_config import get_default_config
        user_api_config = get_default_config(user_id=0)

    prompt = f"""请为PPT页面「{req.page_title}」生成内容要点。

主题：{req.topic}
副标题：{req.subtitle}
行业：{industry_info['name']}
页面类型：{req.page_type}

{full_context}

请生成3-5个bullet points，作为该页的主要内容要点。
每个要点简洁有力，20字以内。
返回JSON格式：{{"content": ["要点1", "要点2", ...]}}"""

    result = call_ai_model(prompt, user_api_config=user_api_config)

    try:
        match = re.search(r'\{.*\}', result, re.DOTALL)
        if match:
            data = json.loads(match.group())
            content = data.get("content", [])
            if content:
                return {"success": True, "content": content}
    except:
        pass

    # 无API Key时的降级预览
    fallback_content = get_fallback_preview(req.topic, req.page_title, req.subtitle)
    return {"success": True, "content": fallback_content}


def get_fallback_preview(topic: str, page_title: str, subtitle: str = "") -> list:
    """无API Key时返回的内容预览降级数据"""
    import hashlib
    # 用topic生成稳定的随机种子
    seed = int(hashlib.md5(f"{topic}{page_title}".encode()).hexdigest()[:8], 16)
    templates = [
        ["市场需求分析", "目标用户画像", "核心竞争力", "商业模式设计", "实施计划"],
        ["问题定义", "解决方案", "技术架构", "开发周期", "预期效果"],
        ["行业背景", "市场机会", "产品功能", "盈利模式", "发展规划"],
        ["用户痛点", "产品价值", "功能模块", "技术方案", "时间规划"],
        ["项目概述", "核心功能", "技术选型", "开发计划", "收益预测"],
    ]
    selected = templates[seed % len(templates)]
    return selected


# ===== SSE进度推送 =====
@app.get("/api/progress/{session_id}")
async def get_progress(session_id: str):
    """SSE端点：推送生成进度"""
    async def event_generator():
        last_status = ""
        while True:
            if session_id in progress_store:
                status = progress_store[session_id]
                if status != last_status:
                    # 发送进度数据
                    yield f"data: {json.dumps(status, ensure_ascii=False)}\n\n"
                    last_status = status
                    # 如果完成了或失败了，发送结束信号
                    if status.get("stage") in ("done", "error"):
                        yield f"data: {json.dumps({'type': 'end'}, ensure_ascii=False)}\n\n"
                        break
            await asyncio.sleep(0.5)
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@app.post("/api/generate_content")
def generate_content(req: ContentRequest):
    """生成PPT内容并渲染 - 专业版（带进度推送）"""
    import traceback
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    # 生成session_id用于追踪进度
    session_id = uuid.uuid4().hex[:8]
    total_pages = len(req.outline)
    progress_store[session_id] = {"stage": "starting", "progress": 0, "message": "正在初始化...", "current_page": 0, "total_pages": total_pages}
    
    def update_progress(stage, progress, message, current_page=0):
        progress_store[session_id] = {
            "stage": stage,
            "progress": progress,
            "message": message,
            "current_page": current_page,
            "total_pages": total_pages
        }
    
    try:
        industry = req.industry or "general"
        update_progress("initializing", 5, "正在初始化PPT文档...")
        
        # 获取用户的模型配置
        user_api_config = None
        if req.model_config_id:
            from model_config import get_config_by_id
            user_api_config = get_config_by_id(req.model_config_id, user_id=0)  # 暂时用0，后续从token获取
            if not user_api_config:
                # 尝试从默认配置获取
                from model_config import get_default_config
                user_api_config = get_default_config(user_id=0)
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 模板颜色配置
        template_colors = {
            "academic": {"title": (26, 54, 93), "accent": (44, 82, 130), "bg": (255, 255, 255)},
            "business": {"title": (44, 82, 130), "accent": (237, 137, 54), "bg": (255, 255, 255)},
            "enterprise": {"title": (30, 58, 95), "accent": (49, 130, 206), "bg": (255, 255, 255)},
            "navy_gold": {"title": (26, 54, 93), "accent": (214, 158, 46), "bg": (255, 255, 255)},
            "simple": {"title": (45, 55, 72), "accent": (226, 232, 240), "bg": (247, 250, 252)},
            "minimal": {"title": (45, 55, 72), "accent": (237, 242, 247), "bg": (255, 255, 255)},
            "gray_elegant": {"title": (45, 55, 72), "accent": (113, 128, 150), "bg": (255, 255, 255)},
            "tech": {"title": (26, 26, 78), "accent": (0, 212, 255), "bg": (255, 255, 255)},
            "cyber": {"title": (15, 15, 35), "accent": (255, 0, 255), "bg": (15, 15, 35)},
            "future": {"title": (26, 5, 51), "accent": (124, 58, 237), "bg": (15, 10, 30)},
            "gradient": {"title": (255, 107, 107), "accent": (78, 205, 196), "bg": (255, 255, 255)},
            "nature": {"title": (39, 103, 73), "accent": (104, 211, 145), "bg": (255, 255, 255)},
            "ocean": {"title": (0, 119, 182), "accent": (0, 180, 216), "bg": (255, 255, 255)},
            "sky": {"title": (116, 185, 255), "accent": (162, 155, 254), "bg": (255, 255, 255)},
            "elegant": {"title": (85, 60, 154), "accent": (159, 122, 234), "bg": (255, 255, 255)},
            "royal": {"title": (76, 29, 149), "accent": (139, 92, 246), "bg": (255, 255, 255)},
            "festive": {"title": (197, 48, 48), "accent": (252, 129, 129), "bg": (255, 255, 255)},
            "classic_blue": {"title": (43, 108, 176), "accent": (99, 179, 237), "bg": (255, 255, 255)},
            "chinese": {"title": (26, 32, 44), "accent": (113, 128, 150), "bg": (255, 255, 255)},
            "ink_wash": {"title": (45, 55, 72), "accent": (160, 174, 192), "bg": (255, 255, 255)},
            "red_gold": {"title": (116, 42, 42), "accent": (214, 158, 46), "bg": (255, 255, 255)},
        }

        colors = template_colors.get(req.template, template_colors["academic"])

        # 如果用户指定了主题色，使用自动配色方案覆盖模板颜色
        if req.brand_color:
            try:
                auto_palette = generate_palette(req.brand_color)
                palette_colors = auto_palette["colors"]
                colors = {
                    "title": hex_to_rgb(palette_colors["primary"]),
                    "accent": hex_to_rgb(palette_colors["accent"]),
                    "bg": hex_to_rgb(palette_colors["background"]),
                    "primary": hex_to_rgb(palette_colors["primary"]),
                    "primary_light": hex_to_rgb(palette_colors["primary_light"]),
                    "primary_dark": hex_to_rgb(palette_colors["primary_dark"]),
                    "surface": hex_to_rgb(palette_colors["surface"]),
                    "text": hex_to_rgb(palette_colors["text_primary"]),
                    "text_secondary": hex_to_rgb(palette_colors["text_secondary"]),
                    "border": hex_to_rgb(palette_colors["border"]),
                }
                update_progress("initializing", 8, "正在应用自定义主题色...")
            except Exception as e:
                print(f"自动配色失败，使用模板默认颜色: {e}")

        # 开始生成幻灯片
        base_progress = 10
        progress_per_page = 80 // total_pages if total_pages > 0 else 80
        
        for idx, page in enumerate(req.outline):
            current_page = idx + 1
            page_progress = base_progress + (current_page - 1) * progress_per_page
            update_progress("generating", int(page_progress), f"正在生成第 {current_page}/{total_pages} 页：{page.get('title', '未命名')[:20]}...", current_page)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page_type = page.get("type", "content")
            chapter_title = page.get("title", "")

            # 为每个章节生成结构化内容
            if page_type == "content":
                update_progress("generating", int(page_progress + 5), f"正在为「{chapter_title[:15]}」生成内容...", current_page)
                kb_context = kb.get_context_for_prompt(chapter_title, industry, top_k=2)
                system_prompt, prompt_template = get_industry_prompt(industry, chapter_title)
                if kb_context:
                    system_prompt = system_prompt + f"\n\n【知识库参考资料】：\n{kb_context}"
                ai_result = call_ai_model(prompt_template, system=system_prompt, user_api_config=user_api_config)
                page_data = parse_ai_content(ai_result)
                page_data["title"] = chapter_title
            elif page_type == "cover":
                page_data = {
                    "title": chapter_title,
                    "subtitle": f"专业PPT解决方案 | {INDUSTRY_TEMPLATES.get(industry, INDUSTRY_TEMPLATES['general'])['name']}",
                    "points": []
                }
            elif page_type == "toc":
                page_data = {
                    "title": "目 录",
                    "subtitle": "CONTENTS",
                    "points": [{"heading": p.get("title", ""), "content": ""} for p in req.outline if p.get("type") == "content"]
                }
            elif page_type == "summary":
                page_data = {
                    "title": "总结与展望",
                    "subtitle": "SUMMARY",
                    "points": [
                        {"heading": "核心要点", "content": "回顾本章核心内容"},
                        {"heading": "实践建议", "content": "如何应用到实际工作中"},
                        {"heading": "未来展望", "content": "发展趋势和机会"}
                    ]
                }
            else:
                page_data = {"title": chapter_title, "points": [{"heading": "内容", "content": "详细说明"}]}

            update_progress("generating", int(page_progress + progress_per_page // 2), f"正在渲染第 {current_page}/{total_pages} 页...", current_page)
            render_structured_slide(slide, page_data, colors, req.template, page_type)

            # 演讲者备注
            if req.speaker_notes and page_type == "content":
                notes_text = generate_speaker_notes(chapter_title, page_data)
                if notes_text and slide.notes_slide:
                    notes_slide = slide.notes_slide
                    notes_tf = notes_slide.notes_text_frame
                    notes_tf.text = notes_text

            # 品牌定制：品牌名/Logo
            if req.brand_name and page_type == "cover":
                add_brand_to_slide(slide, req.brand_name, req.brand_logo, colors)

            # 多语言支持
            if req.language and req.language != "zh":
                page_data = translate_page_content(page_data, req.language)

        # 品牌定制：全局主题色覆盖
        if req.brand_color:
            update_progress("generating", 92, "正在应用全局主题色...")
            apply_brand_color(prs, req.brand_color)

        # 保存
        update_progress("saving", 95, "正在保存PPT文件...")
        filename = f"demo_{uuid.uuid4().hex[:8]}.pptx"
        filepath = OUTPUT_DIR / filename
        prs.save(str(filepath))

        update_progress("done", 100, "PPT生成完成！", total_pages)
        return {"success": True, "download_url": f"/api/download/{filename}", "filename": filename, "session_id": session_id}
        
    except Exception as e:
        error_msg = f"生成失败：{str(e)}"
        print(f"PPT生成错误: {error_msg}\n{traceback.format_exc()}")
        update_progress("error", 0, error_msg, 0)
        return {"success": False, "error": error_msg, "session_id": session_id}


# ===== 演讲者备注生成 =====
def generate_speaker_notes(title: str, page_data: dict) -> str:
    """根据页面内容生成专业演讲者备注"""
    notes = []
    notes.append(f"【章节】{title}")
    notes.append("")
    notes.append("【要点提示】")
    for i, point in enumerate(page_data.get("points", [])[:3], 1):
        heading = point.get("heading", "")
        content = point.get("content", "")
        if heading:
            notes.append(f"• {heading}" + (f"：{content}" if content else ""))
        elif content:
            notes.append(f"• {content}")
    notes.append("")
    notes.append("【演讲建议】")
    notes.append(f"• 开场：可以先回顾上一章节内容，再引入「{title}」")
    notes.append(f"• 重点：强调核心观点，结合实际案例说明")
    notes.append(f"• 过渡：讲解完毕后可适当停顿，询问听众是否有疑问")
    return "\n".join(notes)


# ===== 品牌定制 =====
def add_brand_to_slide(slide, brand_name: str, brand_logo: str, colors: dict):
    """在封面页添加品牌名称和Logo"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    try:
        # 品牌名称 - 左下角
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(3)
        height = Inches(0.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = brand_name
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(*colors["title"])
        p.font.bold = False
        p.alignment = PP_ALIGN.LEFT

        # 如果有logo URL，尝试下载并添加
        if brand_logo and brand_logo.startswith("http"):
            import requests
            from io import BytesIO
            try:
                resp = requests.get(brand_logo, timeout=5)
                if resp.status_code == 200:
                    img_stream = BytesIO(resp.content)
                    # Logo放在右下角
                    left_logo = Inches(10.5)
                    top_logo = Inches(6.6)
                    width_logo = Inches(1.2)
                    height_logo = Inches(0.6)
                    slide.shapes.add_picture(img_stream, left_logo, top_logo, width_logo, height_logo)
            except:
                pass
    except Exception as e:
        print(f"添加品牌信息失败: {e}")


def apply_brand_color(prs, brand_color: str):
    """应用品牌主题色到所有幻灯片（使用自动配色方案）"""
    try:
        auto_palette = generate_palette(brand_color)
        palette_colors = auto_palette["colors"]
        
        # 构建完整的配色方案
        colors = {
            "title": hex_to_rgb(palette_colors["primary"]),
            "accent": hex_to_rgb(palette_colors["accent"]),
            "bg": hex_to_rgb(palette_colors["background"]),
            "primary": hex_to_rgb(palette_colors["primary"]),
            "primary_light": hex_to_rgb(palette_colors["primary_light"]),
            "primary_dark": hex_to_rgb(palette_colors["primary_dark"]),
            "surface": hex_to_rgb(palette_colors["surface"]),
            "text": hex_to_rgb(palette_colors["text_primary"]),
            "text_secondary": hex_to_rgb(palette_colors["text_secondary"]),
            "border": hex_to_rgb(palette_colors["border"]),
        }
        
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
        
        # 遍历所有幻灯片，更新颜色
        for slide in prs.slides:
            for shape in slide.shapes:
                # 更新形状的填充色
                if shape.has_fill and shape.fill.type is not None:
                    try:
                        if shape.fill.type == 1:  # solid fill
                            fill_color = shape.fill.fore_color.rgb
                            # 根据颜色亮度判断替换策略
                            if fill_color:
                                r, g, b = fill_color[0], fill_color[1], fill_color[2]
                                # 替换主色系
                                if _is_similar_color((r,g,b), (26, 54, 93), 50) or \
                                   _is_similar_color((r,g,b), (44, 82, 130), 50) or \
                                   _is_similar_color((r,g,b), (30, 58, 95), 50):
                                    shape.fill.solid()
                                    shape.fill.fore_color.rgb = RGBColor(*colors["title"])
                                # 替换强调色系
                                elif _is_similar_color((r,g,b), (237, 137, 54), 40) or \
                                     _is_similar_color((r,g,b), (44, 82, 130), 40):
                                    shape.fill.solid()
                                    shape.fill.fore_color.rgb = RGBColor(*colors["accent"])
                    except:
                        pass
                
                # 更新文字颜色
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.font:
                            if hasattr(run.font, 'color') and run.font.color and run.font.color.rgb:
                                try:
                                    fc = run.font.color.rgb
                                    r, g, b = fc[0], fc[1], fc[2]
                                    # 替换标题颜色
                                    if _is_similar_color((r,g,b), (26, 54, 93), 40) or \
                                       _is_similar_color((r,g,b), (44, 82, 130), 40):
                                        run.font.color.rgb = RGBColor(*colors["title"])
                                    # 替换强调色
                                    elif _is_similar_color((r,g,b), (237, 137, 54), 30):
                                        run.font.color.rgb = RGBColor(*colors["accent"])
                                except:
                                    pass
    except Exception as e:
        print(f"应用品牌色失败: {e}")


def _is_similar_color(c1: tuple, c2: tuple, threshold: int = 30) -> bool:
    """判断两个颜色是否相似（欧氏距离）"""
    return abs(c1[0] - c2[0])**2 + abs(c1[1] - c2[1])**2 + abs(c1[2] - c2[2])**2 < threshold**2


# ===== 多语言支持 =====
TRANSLATIONS = {
    "en": {
        "目录": "CONTENTS", "总结与展望": "Summary & Outlook", "总结": "Summary",
        "核心要点": "Key Points", "实践建议": "Practical Suggestions", "未来展望": "Future Outlook",
        "专业PPT解决方案": "Professional PPT Solutions",
    },
    "zh-TW": {
        "目录": "目錄", "总结与展望": "總結與展望", "总结": "總結",
        "核心要点": "核心要點", "实践建议": "實踐建議", "未来展望": "未來展望",
        "专业PPT解决方案": "專業PPT解決方案",
    }
}

def translate_page_content(page_data: dict, language: str) -> dict:
    """将页面内容翻译为指定语言"""
    trans = TRANSLATIONS.get(language, {})
    if not trans:
        return page_data
    result = {}
    for key, value in page_data.items():
        if isinstance(value, str):
            result[key] = trans.get(value, value)
        elif isinstance(value, list):
            result[key] = [
                {k: trans.get(v, v) if isinstance(v, str) else v for k, v in item.items()}
                for item in value
            ]
        else:
            result[key] = value
    return result


@app.get("/api/download/{filename}")
async def download(filename: str):
    """下载PPTX"""
    # [SEC-FIX] 防止路径遍历攻击
    if ".." in filename or "/" in filename or "\\" in filename:
        raise HTTPException(status_code=400, detail="无效的文件名")
    filepath = OUTPUT_DIR / filename
    if not filepath.exists():
        raise HTTPException(status_code=404, detail="文件不存在")
    from fastapi.responses import FileResponse
    return FileResponse(str(filepath), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)

@app.get("/api/templates")
async def get_templates():
    """获取模板列表（20款精选模板）"""
    return {"success": True, "templates": [
        {"id": "academic", "name": "学术风", "description": "深蓝/白色，简洁专业", "color": "#1a365d", "accent": "#2c5282"},
        {"id": "business", "name": "商务蓝", "description": "藏蓝/金色，高端大气", "color": "#2c5282", "accent": "#ed8936"},
        {"id": "enterprise", "name": "企业蓝", "description": "稳重企业风格", "color": "#1e3a5f", "accent": "#3182ce"},
        {"id": "navy_gold", "name": "藏金商务", "description": "藏蓝/金色搭配", "color": "#1a365d", "accent": "#d69e2e"},
        {"id": "simple", "name": "简约风", "description": "纯白/浅灰", "color": "#f7fafc", "accent": "#e2e8f0"},
        {"id": "minimal", "name": "极简白", "description": "纯白极简风格", "color": "#ffffff", "accent": "#edf2f7"},
        {"id": "gray_elegant", "name": "灰雅风", "description": "高级灰优雅", "color": "#2d3748", "accent": "#718096"},
        {"id": "tech", "name": "科技风", "description": "渐变蓝/紫色", "color": "#1a1a4e", "accent": "#00d4ff"},
        {"id": "cyber", "name": "赛博朋克", "description": "霓虹科技感", "color": "#0f0f23", "accent": "#ff00ff"},
        {"id": "future", "name": "未来科技", "description": "深空蓝紫色", "color": "#1a0533", "accent": "#7c3aed"},
        {"id": "gradient", "name": "渐变风", "description": "彩色活力", "color": "#ff6b6b", "accent": "#4ecdc4"},
        {"id": "nature", "name": "商务绿", "description": "绿色清新", "color": "#276749", "accent": "#68d391"},
        {"id": "ocean", "name": "海洋蓝", "description": "清爽海风", "color": "#0077b6", "accent": "#00b4d8"},
        {"id": "sky", "name": "天空蓝", "description": "清新天空", "color": "#74b9ff", "accent": "#a29bfe"},
        {"id": "elegant", "name": "商务紫", "description": "紫色高贵", "color": "#553c9a", "accent": "#9f7aea"},
        {"id": "royal", "name": "皇家紫", "description": "皇室贵族风格", "color": "#4c1d95", "accent": "#8b5cf6"},
        {"id": "festive", "name": "商务红", "description": "红色喜庆", "color": "#c53030", "accent": "#fc8181"},
        {"id": "classic_blue", "name": "经典蓝", "description": "经典蓝色", "color": "#2b6cb0", "accent": "#63b3ed"},
        {"id": "chinese", "name": "中国风", "description": "水墨文化", "color": "#1a202c", "accent": "#718096"},
        {"id": "ink_wash", "name": "水墨风", "description": "山水意境", "color": "#2d3748", "accent": "#a0aec0"},
        {"id": "red_gold", "name": "中国红金", "description": "红金传统", "color": "#742a2a", "accent": "#d69e2e"},
    ]}

@app.get("/api/color_palettes")
async def get_color_palettes():
    """获取预设配色方案列表"""
    return {"success": True, "palettes": list_preset_palettes()}

@app.post("/api/generate_palette")
def api_generate_palette(color: str, palette_type: str = "auto"):
    """生成配色方案"""
    try:
        palette = generate_palette(color, palette_type)
        return {"success": True, "palette": palette}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/preview_palette")
async def preview_palette(color: str, palette_type: str = "auto"):
    """预览配色方案 - 生成一张预览图"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        import datetime
        from datetime import datetime as dt
        palette = generate_palette(color, palette_type)
        colors = palette["colors"]
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 全深色背景封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors["background"]))
        bg.line.fill.background()
        
        # 顶部主题色条
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors["primary"]))
        top_bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"配色预览 | {palette['description']}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 颜色展示区域
        color_items = [
            ("primary", "主色", colors["primary"]),
            ("primary_light", "主色浅", colors["primary_light"]),
            ("primary_dark", "主色深", colors["primary_dark"]),
            ("accent", "强调色", colors["accent"]),
            ("accent_light", "强调色浅", colors["accent_light"]),
        ]
        
        for i, (key, label, hex_c) in enumerate(color_items):
            x = 0.5 + i * 2.5
            rgb = hex_to_rgb(hex_c)
            is_light = rgb[0]*0.299 + rgb[1]*0.587 + rgb[2]*0.114 > 180
            
            color_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.2), Inches(1.8))
            color_box.fill.solid()
            color_box.fill.fore_color.rgb = RGBColor(*rgb)
            color_box.line.color.rgb = RGBColor(200, 200, 200)
            
            txt_box = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(2.2), Inches(0.5))
            tf = txt_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_to_rgb(colors["text_primary"]))
            p.alignment = 1  # center
            
            code_box = slide.shapes.add_textbox(Inches(x), Inches(3.9), Inches(2.2), Inches(0.4))
            tf = code_box.text_frame
            p = tf.paragraphs[0]
            p.text = hex_c.upper()
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(*hex_to_rgb(colors["text_secondary"]))
            p.alignment = 1
        
        # 辅助色展示
        aux_items = [
            ("background", "背景", colors["background"]),
            ("surface", "卡片", colors["surface"]),
            ("text_primary", "主文字", colors["text_primary"]),
            ("text_secondary", "次文字", colors["text_secondary"]),
            ("border", "边框", colors["border"]),
        ]
        
        for i, (key, label, hex_c) in enumerate(aux_items):
            x = 0.5 + i * 2.5
            rgb = hex_to_rgb(hex_c)
            
            color_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.6), Inches(2.2), Inches(1.2))
            color_box.fill.solid()
            color_box.fill.fore_color.rgb = RGBColor(*rgb)
            color_box.line.color.rgb = RGBColor(200, 200, 200)
            
            txt_box = slide.shapes.add_textbox(Inches(x), Inches(5.9), Inches(2.2), Inches(0.4))
            tf = txt_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label} {hex_c.upper()}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*hex_to_rgb(colors["text_primary"]))
            p.alignment = 1
        
        # 底部示例文字
        example_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = example_box.text_frame
        p = tf.paragraphs[0]
        p.text = "演示示例文字 | 1234567890 ABCabc"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*hex_to_rgb(colors["text_primary"]))
        
        filename = f"palette_preview_{dt.now().strftime('%Y%m%d%H%M%S')}.pptx"
        filepath = OUTPUT_DIR / filename
        prs.save(str(filepath))
        
        return {
            "success": True,
            "download_url": f"/api/download/{filename}",
            "filename": filename,
            "palette": palette
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def hex_to_rgb(hex_color: str) -> tuple:
    """HEX转RGB"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

@app.get("/api/industries")
async def get_industries():
    """获取行业列表"""
    return {"success": True, "industries": [
        {"id": "education", "name": "教育培训", "icon": "🎓"},
        {"id": "medical", "name": "医疗健康", "icon": "🏥"},
        {"id": "ecommerce", "name": "电商零售", "icon": "🛒"},
        {"id": "finance", "name": "金融投资", "icon": "💰"},
        {"id": "technology", "name": "科技互联网", "icon": "💻"},
        {"id": "government", "name": "政府企业", "icon": "🏛️"},
        {"id": "realestate", "name": "房产建筑", "icon": "🏗️"},
        {"id": "media", "name": "传媒广告", "icon": "📺"},
        {"id": "manufacture", "name": "制造业", "icon": "⚙️"},
        {"id": "general", "name": "通用场景", "icon": "📋"},
    ]}

@app.post("/api/convert_document")
async def convert_document(req: DocumentConvertRequest):
    """文档一键转PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    industry_config = INDUSTRY_TEMPLATES.get(req.industry, INDUSTRY_TEMPLATES["general"])
    template_id = req.template if req.template in [t for i in INDUSTRY_TEMPLATES.values() for t in i["templates"]] else industry_config["templates"][0]

    # 获取用户配置的模型
    user_api_config = None
    if req.model_config_id:
        from model_config import get_config_by_id, get_default_config
        user_api_config = get_config_by_id(req.model_config_id, user_id=0)
        if not user_api_config:
            user_api_config = get_default_config(user_id=0)
    else:
        from model_config import get_default_config
        user_api_config = get_default_config(user_id=0)

    template_colors = {
        "academic": {"title": (26, 54, 93), "accent": (44, 82, 130)},
        "business": {"title": (44, 82, 130), "accent": (237, 137, 54)},
        "enterprise": {"title": (30, 58, 95), "accent": (49, 130, 206)},
        "navy_gold": {"title": (26, 54, 93), "accent": (214, 158, 46)},
        "simple": {"title": (45, 55, 72), "accent": (226, 232, 240)},
        "nature": {"title": (39, 103, 73), "accent": (104, 211, 145)},
        "tech": {"title": (26, 26, 78), "accent": (0, 212, 255)},
        "cyber": {"title": (15, 15, 35), "accent": (255, 0, 255)},
        "gradient": {"title": (255, 107, 107), "accent": (78, 205, 196)},
        "ocean": {"title": (0, 119, 182), "accent": (0, 180, 216)},
        "elegant": {"title": (85, 60, 154), "accent": (159, 122, 234)},
        "royal": {"title": (76, 29, 149), "accent": (139, 92, 246)},
        "festive": {"title": (197, 48, 48), "accent": (252, 129, 129)},
        "classic_blue": {"title": (43, 108, 176), "accent": (99, 179, 237)},
        "gray_elegant": {"title": (45, 55, 72), "accent": (113, 128, 150)},
        "ink_wash": {"title": (45, 55, 72), "accent": (160, 174, 192)},
    }
    colors = template_colors.get(template_id, template_colors["business"])

    # 用AI分析文档内容，生成大纲
    system_prompt, _ = get_industry_prompt(req.industry, "文档分析")
    prompt = f"""请分析以下文档内容，生成专业的PPT大纲结构。
文档内容：
{req.text[:3000]}

要求：
- 5-8页幻灯片
- 包含封面、目录、内容章节、总结
- 返回JSON格式：{{"pages": [{{"title": "标题", "type": "cover/toc/content/summary"}}]}}"""

    ai_result = call_ai_model(prompt, system=system_prompt, user_api_config=user_api_config)

    outline = []
    try:
        match = re.search(r'\{.*\}', ai_result, re.DOTALL)
        if match:
            data = json.loads(match.group())
            outline = data.get("pages", [])
    except:
        pass

    if not outline:
        lines = req.text.split('\n')
        outline = [
            {"title": "文档标题", "type": "cover"},
            {"title": "目录", "type": "toc"},
        ]
        for i, line in enumerate(lines[:6]):
            if line.strip():
                outline.append({"title": line.strip()[:50], "type": "content"})
        outline.append({"title": "总结", "type": "summary"})

    # 生成PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for page in outline:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_type = page.get("type", "content")
        chapter_title = page.get("title", "")

        if page_type == "content":
            system_prompt, prompt_template = get_industry_prompt(req.industry, chapter_title)
            ai_result = call_ai_model(prompt_template, system=system_prompt, user_api_config=user_api_config)
            page_data = parse_ai_content(ai_result)
            page_data["title"] = chapter_title
        elif page_type == "cover":
            page_data = {"title": chapter_title, "subtitle": f"文档转换 | {industry_config['name']}", "points": []}
        elif page_type == "toc":
            page_data = {"title": "目 录", "subtitle": "CONTENTS", "points": [{"heading": p.get("title", ""), "content": ""} for p in outline if p.get("type") == "content"]}
        elif page_type == "summary":
            page_data = {"title": "总结", "subtitle": "SUMMARY", "points": [{"heading": "核心要点", "content": "回顾全文核心内容"}]}
        else:
            page_data = {"title": chapter_title, "points": [{"heading": "内容", "content": "详细说明"}]}

        render_structured_slide(slide, page_data, colors, template_id, page_type)

    filename = f"doc_{uuid.uuid4().hex[:8]}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))

    return {
        "success": True,
        "download_url": f"/api/download/{filename}",
        "filename": filename,
        "industry": industry_config["name"],
        "template": template_id,
        "outline": outline
    }

@app.post("/api/upload_document")
async def upload_document(file: bytes = None, industry: str = "general", template: str = "business"):
    """上传文档文件并转换为PPT"""
    if not file:
        raise HTTPException(status_code=400, detail="请上传文件")

    text_content = ""
    filename = f"upload_{uuid.uuid4().hex[:8]}"

    try:
        text_content = file.decode('utf-8', errors='ignore')
    except:
        text_content = str(file)

    if len(text_content) < 50:
        return {"success": False, "message": "文件内容读取失败，请上传txt格式或直接粘贴文本内容"}

    req = DocumentConvertRequest(text=text_content, industry=industry, template=template)
    return await convert_document(req)

@app.post("/api/digital_human")
async def digital_human(req: DigitalHumanRequest):
    """数字人播报生成"""
    try:
        avatars = {
            "default": {"voice": "zh-CN-Female", "style": "professional"},
            "female": {"voice": "zh-CN-Female", "style": "warm"},
            "male": {"voice": "zh-CN-Male", "style": "authoritative"},
            "cartoon": {"voice": "zh-CN-Child", "style": "friendly"}
        }
        avatar_config = avatars.get(req.avatar, avatars["default"])

        return {
            "success": True,
            "message": "数字人功能开发中",
            "avatar": req.avatar,
            "avatar_config": avatar_config,
            "note": "数字人功能需要配置API后使用"
        }
    except Exception as e:
        return {"success": False, "message": str(e)}

class DocumentImportRequest(BaseModel):
    """文档导入请求 - 粘贴内容直接生成PPT"""
    content: str  # 原始文本/Markdown内容
    topic: Optional[str] = ""  # 可选主题，用于AI理解上下文
    template: str = "academic"
    industry: Optional[str] = "general"
    language: Optional[str] = "zh"

@app.post("/api/import_document")
async def import_document(req: DocumentImportRequest):
    """文档导入 - 粘贴文本/Markdown，直接解析成PPT大纲并生成"""
    from datetime import datetime
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # 1. 用AI分析文档结构，提取标题和内容要点
        analysis_prompt = f"""你是一个专业的PPT内容规划师。请分析下面的文档内容，提取出适合制作PPT的结构。

要求：
1. 识别文档的主要主题和章节结构
2. 每个章节提炼出3-5个核心要点
3. 识别文档中的数据、图表需求（如果有）
4. 判断适合的PPT类型（封面/目录/内容页/总结页）

输出格式（JSON）：
{{
  "title": "PPT主题",
  "cover_subtitle": "副标题",
  "outline": [
    {{"title": "章节标题", "type": "cover/toc/content/summary", "points": ["要点1", "要点2"], "has_chart": false, "chart_type": null}},
    ...
  ]
}}

文档内容：
{req.content[:3000]}"""

        industry = req.industry or "general"
        system_prompt = "你是一个专业的PPT内容规划专家，擅长将文档内容结构化并生成适合PPT展示的大纲。始终返回JSON格式。"
        
        ai_raw = call_ai_model(analysis_prompt, system=system_prompt)
        
        # 解析AI返回
        import re, json as json_lib
        json_match = re.search(r'\{[\s\S]*\}', ai_raw)
        if json_match:
            parsed = json_lib.loads(json_match.group())
        else:
            # fallback: 简单按段落分割
            lines = [l.strip() for l in req.content.split('\n') if l.strip()]
            parsed = {
                "title": req.topic or "文档导入PPT",
                "cover_subtitle": "",
                "outline": [{"title": lines[i] if i < len(lines) else f"第{i}点", "type": "content", "points": [], "has_chart": False} for i in range(min(len(lines), 6))]
            }

        # 2. 生成PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        template_colors = {
            "academic": {"title": (26, 54, 93), "accent": (44, 82, 130), "bg": (255, 255, 255)},
            "business": {"title": (44, 82, 130), "accent": (237, 137, 54), "bg": (255, 255, 255)},
            "tech": {"title": (26, 26, 78), "accent": (0, 212, 255), "bg": (255, 255, 255)},
        }
        colors = template_colors.get(req.template, template_colors["academic"])

        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = parsed.get("title", req.topic or "文档导入PPT")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors["title"])
        if parsed.get("cover_subtitle"):
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            sub_box.text_frame.paragraphs[0].text = parsed["cover_subtitle"]
            sub_box.text_frame.paragraphs[0].font.size = Pt(24)
            sub_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["accent"])

        # 目录
        outline = parsed.get("outline", [])
        if len(outline) > 1:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
            title_box.text_frame.paragraphs[0].text = "目录"
            title_box.text_frame.paragraphs[0].font.size = Pt(32)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["title"])

            toc_items = [item for item in outline if item.get("type") != "cover"]
            for i, item in enumerate(toc_items[:8]):
                y = Inches(1.8 + i * 0.65)
                num_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(0.5), Inches(0.5))
                num_box.text_frame.paragraphs[0].text = f"{i+1}"
                num_box.text_frame.paragraphs[0].font.size = Pt(20)
                num_box.text_frame.paragraphs[0].font.bold = True
                num_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["accent"])
                txt_box = slide.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(0.5))
                txt_box.text_frame.paragraphs[0].text = item.get("title", "")
                txt_box.text_frame.paragraphs[0].font.size = Pt(18)

        # 内容页
        for item in outline:
            if item.get("type") in ("cover", "toc"):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9))
            title_box.text_frame.paragraphs[0].text = item.get("title", "")
            title_box.text_frame.paragraphs[0].font.size = Pt(28)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["title"])
            # 要点
            points = item.get("points", [])
            if points:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                for j, point in enumerate(points):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)

        # 保存
        filename = f"doc_import_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
        filepath = OUTPUT_DIR / filename
        prs.save(str(filepath))

        return {
            "success": True,
            "file": filename,
            "download_url": f"/api/download/{filename}",
            "outline": parsed.get("outline", []),
            "title": parsed.get("title", req.topic or "文档导入PPT"),
            "message": "文档导入成功"
        }
    except Exception as e:
        return {"success": False, "message": f"文档导入失败: {str(e)}"}


class ContentPolishRequest(BaseModel):
    """AI内容润色请求"""
    content: str  # 要润色的文本
    mode: str = "polish"  # polish(润色)/expand(扩写)/simplify(精简)
    language: Optional[str] = "zh"
    model_config_id: Optional[int] = None  # 用户配置的模型ID

@app.post("/api/polish_content")
async def polish_content(req: ContentPolishRequest):
    """AI内容润色 - 扩写/精简/改写"""
    try:
        # 获取用户配置的模型
        user_api_config = None
        if req.model_config_id:
            from model_config import get_config_by_id, get_default_config
            user_api_config = get_config_by_id(req.model_config_id, user_id=0)
            if not user_api_config:
                user_api_config = get_default_config(user_id=0)
        else:
            from model_config import get_default_config
            user_api_config = get_default_config(user_id=0)

        mode_descriptions = {
            "polish": "语言更流畅、专业、有说服力",
            "expand": "在原有内容基础上适当扩展，加入更多细节和案例",
            "simplify": "简化内容，保留核心要点，表达更简洁清晰"
        }
        desc = mode_descriptions.get(req.mode, "优化语言表达")
        
        prompt = f"""请将以下内容进行{desc}：

原文：
{req.content}

要求：
1. 保持原文的核心意思不变
2. 润色后的内容要专业、流畅、有逻辑
3. 适合用于PPT演示文稿
4. 如果是扩写模式，请添加相关的案例或数据支持

请直接输出润色后的内容，不要加额外说明。"""
        
        result = call_ai_model(prompt, system="你是一个专业的中文内容编辑，擅长PPT内容优化。", user_api_config=user_api_config)
        
        return {
            "success": True,
            "original": req.content,
            "polished": result,
            "mode": req.mode,
            "message": "内容润色成功"
        }
    except Exception as e:
        return {"success": False, "message": f"润色失败: {str(e)}"}


class ImageSearchRequest(BaseModel):
    """智能配图请求"""
    keyword: str  # 配图关键词
    style: Optional[str] = "photo"  # photo/illustration/vector

def get_wiki_image(keyword: str, cache_dir: str = "/tmp/ppt_images") -> dict:
    """从Wikipedia/Wikimedia获取配图（免费无key，可商用）- 网络不可用时快速降级"""
    import os, urllib.request, json, urllib.parse, socket
    
    os.makedirs(cache_dir, exist_ok=True)
    safe_kw = keyword.replace("/", "-").replace(" ", "_")[:50]
    img_cache = f"{cache_dir}/{safe_kw}_wiki.jpg"
    
    # 设置短超时，避免网络不可用时卡住
    socket_timeout = 3
    
    try:
        # 第一步：搜索Wikipedia找到相关页面
        search_url = f"https://en.wikipedia.org/w/api.php?action=opensearch&search={urllib.parse.quote(keyword)}&limit=1&namespace=0&format=json"
        
        req = urllib.request.Request(search_url, headers={"User-Agent": "DemoPPT/1.0"})
        socket.setdefaulttimeout(socket_timeout)
        with urllib.request.urlopen(req, timeout=socket_timeout) as resp:
            search_data = json.loads(resp.read())
        
        if len(search_data) < 2 or not search_data[1]:
            raise ValueError("No search results")
        
        page_title = search_data[1][0]
    except:
        page_title = keyword  # 降级使用原始关键词
    
    # 第二步：获取该页面的图片
    wiki_url = f"https://en.wikipedia.org/w/api.php?action=query&titles={urllib.parse.quote(page_title)}&prop=pageimages&format=json&pithumbsize=800&redirects=1"
    
    try:
        req = urllib.request.Request(wiki_url, headers={"User-Agent": "DemoPPT/1.0"})
        with urllib.request.urlopen(req, timeout=socket_timeout) as resp:
            data = json.loads(resp.read())
        
        pages = data.get("query", {}).get("pages", {})
        for page_id, page_data in pages.items():
            if page_id != "-1" and "thumbnail" in page_data:
                img_url = page_data["thumbnail"]["source"]
                
                try:
                    req2 = urllib.request.Request(img_url, headers={"User-Agent": "Mozilla/5.0"})
                    with urllib.request.urlopen(req2, timeout=socket_timeout) as r, open(img_cache, "wb") as f:
                        f.write(r.read())
                    
                    return {
                        "success": True,
                        "keyword": keyword,
                        "images": [{
                            "url": f"file://{img_cache}",
                            "local_path": img_cache,
                            "thumb": f"file://{img_cache}",
                            "credit": f"Wikipedia: {page_data.get('title', page_title)}",
                            "license": "CC BY-SA/Wikimedia Commons"
                        }]
                    }
                except:
                    pass
    except:
        pass
    
    # 网络不可用时返回占位符
    return {
        "success": True,
        "keyword": keyword,
        "images": [{
            "url": "",
            "local_path": "",
            "thumb": "",
            "credit": "（网络不可用，图片功能暂不可用）",
            "license": "N/A"
        }]
    }

@app.get("/api/search_image")
def search_image(keyword: str, style: str = "photo"):
    """智能配图 - Wikipedia/Wikimedia免费图片（可商用）"""
    import urllib.parse
    return get_wiki_image(keyword)


class VideoExportRequest(BaseModel):
    """视频导出请求"""
    pptx_path: str  # PPT文件路径
    duration_per_slide: float = 3.0  # 每页停留秒数
    transition: str = "fade"  # 切换效果 fade/slide

@app.post("/api/export_video")
async def export_video(req: VideoExportRequest):
    """PPT转视频导出 - 将PPTX转换为MP4视频"""
    try:
        import subprocess
        import tempfile
        import os
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # 检查ffmpeg
        try:
            subprocess.run(["ffmpeg", "-version"], capture_output=True, check=True)
        except:
            return {"success": False, "message": "视频导出功能需要安装ffmpeg，当前环境未安装"}
        
        pptx_path = req.pptx_path
        if not os.path.exists(pptx_path):
            return {"success": False, "message": "PPTX文件不存在"}
        
        # 创建临时目录
        temp_dir = tempfile.mkdtemp()
        frames_dir = os.path.join(temp_dir, "frames")
        os.makedirs(frames_dir, exist_ok=True)
        
        # 用matplotlib渲染每页幻灯片为图片
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import matplotlib.patches as patches
            import numpy as np
            
            prs = Presentation(pptx_path)
            total = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                fig, ax = plt.subplots(figsize=(16, 9), dpi=150)
                ax.set_xlim(0, 16)
                ax.set_ylim(0, 9)
                ax.axis('off')
                fig.patch.set_facecolor('white')
                
                # 提取文字
                title_text = ""
                content_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        text = shape.text_frame.text.strip()
                        if text:
                            if not title_text:
                                title_text = text[:100]
                            else:
                                content_texts.append(text[:200])
                
                # 绘制标题
                if title_text:
                    ax.text(8, 7.5, title_text, fontsize=32, fontweight='bold', 
                            ha='center', va='center', color='#1a365d')
                
                # 绘制内容
                y_pos = 6.0
                for j, txt in enumerate(content_texts[:5]):
                    ax.text(1.5, y_pos - j*0.8, f"• {txt}", fontsize=14,
                            ha='left', va='top', color='#2d3748', wrap=True)
                
                # 页码
                ax.text(15.5, 0.5, f"{i+1}/{total}", fontsize=12, 
                        ha='right', va='bottom', color='#718096')
                
                # 保存帧
                frame_path = os.path.join(frames_dir, f"frame_{i:04d}.png")
                plt.savefig(frame_path, bbox_inches='tight', facecolor='white')
                plt.close()
            
            plt.close('all')
            
        except Exception as e:
            return {"success": False, "message": f"幻灯片渲染失败: {str(e)}"}
        
        # 用ffmpeg合成视频
        output_filename = f"video_{datetime.now().strftime('%Y%m%d%H%M%S')}.mp4"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        # 计算参数
        fps = 1 / req.duration_per_slide
        total_duration = total * req.duration_per_slide
        
        # ffmpeg命令
        cmd = [
            "ffmpeg", "-y",
            "-framerate", str(fps),
            "-i", os.path.join(frames_dir, "frame_%04d.png"),
            "-c:v", "libx264",
            "-pix_fmt", "yuv420p",
            "-t", str(total_duration),
            "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2",
            output_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        # 清理临时文件
        import shutil
        shutil.rmtree(temp_dir, ignore_errors=True)
        
        if result.returncode != 0:
            return {"success": False, "message": f"视频合成失败: {result.stderr[-200:]}"}
        
        return {
            "success": True,
            "file": output_filename,
            "download_url": f"/api/download/{output_filename}",
            "duration": total_duration,
            "slides": total,
            "message": "视频导出成功"
        }
        
    except Exception as e:
        return {"success": False, "message": f"视频导出失败: {str(e)}"}


class BeautifyRequest(BaseModel):
    """一键美化请求"""
    content: str  # 要美化的内容
    style: str = "modern"  # modern/elegant/vivid/tech
    model_config_id: Optional[int] = None  # 用户配置的模型ID

@app.post("/api/beautify")
async def beautify_content(req: BeautifyRequest):
    """一键美化 - AI自动优化PPT内容排版和表达"""
    try:
        # 获取用户配置的模型
        user_api_config = None
        if req.model_config_id:
            from model_config import get_config_by_id, get_default_config
            user_api_config = get_config_by_id(req.model_config_id, user_id=0)
            if not user_api_config:
                user_api_config = get_default_config(user_id=0)
        else:
            from model_config import get_default_config
            user_api_config = get_default_config(user_id=0)

        style_descriptions = {
            "modern": "现代简约风格，语言简洁有力，适合商业演示",
            "elegant": "优雅专业风格，语言精致考究，适合正式场合",
            "vivid": "生动活泼风格，语言有趣有感染力，适合教育培训",
            "tech": "科技感风格，语言专业前沿，适合技术分享"
        }
        desc = style_descriptions.get(req.style, "现代专业风格")
        
        prompt = f"""你是一个专业的PPT内容美化专家。请将以下内容进行一键美化，使其更适合PPT演示。

风格要求：{desc}

原文：
{req.content}

请直接输出美化后的内容，保持JSON格式：
{{
  "title": "优化后的标题",
  "subtitle": "副标题（可选）",
  "points": ["要点1", "要点2", "要点3"],
  "takeaway": "核心结论/金句"
}}

要求：
1. 标题要简洁有力，吸引注意力
2. 每个要点控制在20字以内，观点鲜明
3. 核心结论要有记忆点
4. 适合演讲演示使用"""
        
        result = call_ai_model(prompt, system="你是一个专业的PPT内容美化专家，始终返回JSON格式。", user_api_config=user_api_config)
        
        import re, json as json_lib
        json_match = re.search(r'\{[\s\S]*\}', result)
        if json_match:
            parsed = json_lib.loads(json_match.group())
            return {
                "success": True,
                "original": req.content,
                "beautified": parsed,
                "style": req.style,
                "message": "一键美化成功"
            }
        else:
            return {"success": False, "message": "美化解析失败，请重试"}
            
    except Exception as e:
        return {"success": False, "message": f"美化失败: {str(e)}"}


# ============================================================
# 知识库 API
# ============================================================

class KBDocumentRequest(BaseModel):
    content: str
    filename: str
    file_type: str = "txt"
    title: str = ""


@app.get("/api/kb/stats")
def kb_stats():
    """知识库统计"""
    return kb.get_stats()


@app.get("/api/kb/documents")
def kb_list_documents():
    """列出知识库中所有文档"""
    return {"success": True, "documents": kb.list_documents()}


@app.post("/api/kb/documents")
def kb_add_document(req: KBDocumentRequest):
    """添加文档到知识库"""
    try:
        doc_id = kb.add_document(
            content=req.content,
            filename=req.filename,
            file_type=req.file_type,
            title=req.title
        )
        return {"success": True, "doc_id": doc_id, "message": "文档已添加到知识库"}
    except Exception as e:
        return {"success": False, "message": f"添加失败: {str(e)}"}


@app.post("/api/kb/upload")
async def kb_upload():
    """上传文件到知识库"""
    from fastapi import UploadFile, File
    
    try:
        upload_file = await UploadFile.create(File(...))
        content = await upload_file.read()
        
        # 解析文档
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=upload_file.filename) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        
        text = parse_document(tmp_path)
        os.unlink(tmp_path)
        
        if text.startswith("["):
            return {"success": False, "message": text}
        
        doc_id = kb.add_document(
            content=text,
            filename=upload_file.filename,
            file_type=os.path.splitext(upload_file.filename)[1].lower()
        )
        
        return {"success": True, "doc_id": doc_id, "filename": upload_file.filename, "message": "上传成功"}
    except Exception as e:
        return {"success": False, "message": f"上传失败: {str(e)}"}


@app.delete("/api/kb/documents/{doc_id}")
def kb_delete_document(doc_id: str):
    """删除知识库文档"""
    success = kb.delete_document(doc_id)
    return {"success": success, "message": "删除成功" if success else "删除失败"}


@app.get("/api/kb/search")
def kb_search(query: str, top_k: int = 5, industry: str = ""):
    """搜索知识库"""
    results = kb.search(query, top_k=top_k)
    return {"success": True, "query": query, "results": results, "count": len(results)}


@app.get("/api/kb/context")
def kb_get_context(topic: str, industry: str = "", top_k: int = 5):
    """获取用于PPT生成的上下文（知识库召回）"""
    context = kb.get_context_for_prompt(topic, industry, top_k)
    return {"success": True, "context": context, "has_content": bool(context)}


@app.get("/api/search_industry")
def api_search_industry(q: str, industry: str = ""):
    """联网搜索行业数据"""
    result = search_and_get_context(q, industry)
    return {"success": True, **result}


@app.get("/api/industry_kb")
def api_industry_kb(industry: str):
    """获取预置行业知识库内容"""
    # 兼容旧key名称
    KEY_MAP = {"ecommerce": "retail", "medical": "healthcare"}
    industry_key = KEY_MAP.get(industry, industry)
    kb_data = INDUSTRY_KB.get(industry_key, INDUSTRY_KB.get("general", {}))
    return {"success": True, "industry": industry, "industry_key": industry_key, "data": kb_data}


HIST_DB = Path(__file__).parent / "history.db"

def _init_history_db():
    import sqlite3
    conn = sqlite3.connect(str(HIST_DB))
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS ppt_history (
        id TEXT PRIMARY KEY, topic TEXT, subtitle TEXT, industry TEXT,
        template TEXT, outline_json TEXT, filename TEXT, file_size INTEGER,
        created_at TEXT, download_url TEXT)""")
    conn.commit()
    conn.close()

_init_history_db()

class HistoryRecord(BaseModel):
    topic: str
    subtitle: Optional[str] = ""
    industry: Optional[str] = ""
    template: Optional[str] = ""
    outline_json: Optional[str] = ""
    filename: str
    file_size: int = 0
    download_url: Optional[str] = ""

@app.post("/api/history")
def add_history(req: HistoryRecord):
    """保存PPT生成记录"""
    import sqlite3
    from datetime import datetime
    import uuid
    hist_id = str(uuid.uuid4())[:12]
    now = datetime.now().isoformat()
    try:
        conn = sqlite3.connect(str(HIST_DB))
        c = conn.cursor()
        c.execute("""INSERT INTO ppt_history (id, topic, subtitle, industry, template, outline_json, filename, file_size, created_at, download_url)
                     VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (hist_id, req.topic, req.subtitle, req.industry, req.template, req.outline_json, req.filename, req.file_size, now, req.download_url))
        conn.commit()
        conn.close()
        return {"success": True, "id": hist_id}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/api/history")
def list_history(limit: int = 50):
    """获取历史记录列表"""
    import sqlite3
    try:
        conn = sqlite3.connect(str(HIST_DB))
        c = conn.cursor()
        c.execute("SELECT id, topic, subtitle, industry, template, filename, file_size, created_at, download_url FROM ppt_history ORDER BY created_at DESC LIMIT ?", (limit,))
        rows = c.fetchall()
        conn.close()
        records = [{"id": r[0], "topic": r[1], "subtitle": r[2], "industry": r[3], "template": r[4], "filename": r[5], "file_size": r[6], "date": r[7][:16], "url": r[8]} for r in rows]
        return {"success": True, "records": records}
    except Exception as e:
        return {"success": False, "error": str(e), "records": []}

@app.delete("/api/history/{hist_id}")
def delete_history(hist_id: str):
    """删除历史记录"""
    import sqlite3
    try:
        conn = sqlite3.connect(str(HIST_DB))
        c = conn.cursor()
        c.execute("DELETE FROM ppt_history WHERE id = ?", (hist_id,))
        conn.commit()
        conn.close()
        return {"success": True}
    except:
        return {"success": False}

# ============ 认证相关路由 ============
from user_auth import register, login, logout, verify_token, get_user_info, update_user_info, change_password, send_sms_code, verify_sms_code
from subscription import get_all_plans, get_plan, create_subscription, pay_order, get_user_subscription, check_user_access, cancel_subscription, get_user_orders
from model_config import add_config, get_configs, get_default_config, get_config_by_id, update_config, delete_config, set_default, call_model

@app.get("/api/auth/me")
def get_me(authorization: str = ""):
    """获取当前登录用户信息"""
    if not authorization.startswith("Bearer "):
        return {"success": False, "error": "未登录"}
    token = authorization[7:]
    user = verify_token(token)
    if not user:
        return {"success": False, "error": "Token无效"}
    
    sub = get_user_subscription(user["user_id"])
    return {
        "success": True,
        "user": {
            "id": user["user_id"],
            "phone": user["phone"],
            "nickname": user["nickname"],
            "avatar_url": user["avatar_url"]
        },
        "subscription": sub
    }

from pydantic import BaseModel

class RegisterReq(BaseModel):
    phone: str
    password: str
    nickname: str = ""

class LoginReq(BaseModel):
    phone: str
    password: str
    device_info: str = ""

class SendSmsReq(BaseModel):
    phone: str

class PlanCodeReq(BaseModel):
    plan_code: str
    payment_method: str = "wechat"

class VerifySmsReq(BaseModel):
    phone: str
    code: str

class ChangePasswordReq(BaseModel):
    old_password: str
    new_password: str

class AddModelReq(BaseModel):
    name: str
    api_base: str
    api_key: str
    model_name: str
    provider: str = "custom"
    is_default: int = 0

class UpdateModelReq(BaseModel):
    name: str = None
    api_base: str = None
    api_key: str = None
    model_name: str = None
    provider: str = None
    is_default: int = None
    status: int = None

def get_user_from_header(authorization: str = Header(None)) -> dict | None:
    """从Authorization header提取并验证token，返回用户信息或None"""
    if not authorization:
        return None
    if not authorization.startswith("Bearer "):
        return None
    token = authorization[7:]
    return verify_token(token)

@app.post("/api/auth/register")
def api_register(req: RegisterReq):
    """用户注册"""
    result = register(req.phone, req.password, req.nickname)
    return result

@app.post("/api/auth/login")
def api_login(req: LoginReq):
    """用户登录"""
    result = login(req.phone, req.password, req.device_info)
    return result

@app.post("/api/auth/logout")
def api_logout(user = Depends(get_user_from_header)):
    """用户登出"""
    if not user:
        return {"success": False, "error": "未登录"}
    token = None  # 登出不需要token，前端已删除localStorage
    return {"success": True, "message": "已退出登录"}

@app.post("/api/auth/send_sms")
def api_send_sms(phone: str):
    """发送短信验证码"""
    import secrets
    code = str(secrets.randbelow(9000) + 1000)  # 4位验证码
    result = send_sms_code(phone, code)
    return result

@app.post("/api/auth/verify_sms")
def api_verify_sms(phone: str, code: str):
    """验证短信验证码"""
    success = verify_sms_code(phone, code)
    if success:
        # 验证成功，可自动登录或返回token
        return {"success": True, "message": "验证成功"}
    return {"success": False, "error": "验证码错误或已过期"}

@app.put("/api/auth/nickname")
def api_update_nickname(nickname: str, user = Depends(get_user_from_header)):
    """修改昵称"""
    if not user:
        return {"success": False, "error": "未登录"}
    return update_user_info(user["user_id"], nickname=nickname)

@app.put("/api/auth/password")
def api_change_password(req: ChangePasswordReq, user = Depends(get_user_from_header)):
    """修改密码"""
    if not user:
        return {"success": False, "error": "未登录"}
    return change_password(user["user_id"], req.old_password, req.new_password)

# ============ 订阅相关路由 ============
@app.get("/api/plans")
def api_get_plans():
    """获取所有套餐"""
    plans = get_all_plans()
    return {"success": True, "plans": plans}

@app.get("/api/subscription")
def api_get_subscription(user = Depends(get_user_from_header)):
    """获取当前订阅"""
    if not user:
        return {"success": False, "error": "未登录"}
    sub = get_user_subscription(user["user_id"])
    return {"success": True, "subscription": sub}

@app.post("/api/subscription/create")
def api_create_subscription(req: PlanCodeReq, user = Depends(get_user_from_header)):
    """创建订阅订单"""
    if not user:
        return {"success": False, "error": "未登录"}
    return create_subscription(user["user_id"], req.plan_code, req.payment_method)

@app.post("/api/subscription/cancel")
def api_cancel_subscription(user = Depends(get_user_from_header)):
    """取消订阅（关闭自动续费）"""
    if not user:
        return {"success": False, "error": "未登录"}
    return cancel_subscription(user["user_id"])

@app.get("/api/subscription/check")
def api_check_access(user = Depends(get_user_from_header)):
    """检查访问权限"""
    if not user:
        return {"success": False, "has_access": False, "error": "未登录"}
    return check_user_access(user["user_id"])

@app.get("/api/orders")
def api_get_orders(user = Depends(get_user_from_header)):
    """获取订单列表"""
    if not user:
        return {"success": False, "error": "未登录"}
    orders = get_user_orders(user["user_id"])
    return {"success": True, "orders": orders}


# ============ 模型配置路由 ============
@app.get("/api/models")
def api_get_models(user = Depends(get_user_from_header)):
    """获取模型列表（预置+用户配置）"""
    # 预置模型列表
    models_list = []
    for key, config in AI_MODELS.items():
        has_key = bool(os.environ.get(config["env_key"], ""))
        models_list.append({
            "id": key,
            "name": config["name"],
            "provider": config["provider"],
            "models": config["models"],
            "default_model": config["default_model"],
            "has_key": has_key,
        })

    active = os.environ.get("ACTIVE_MODEL", DEFAULT_MODEL)

    # 用户自定义配置（需要登录）
    user_configs = []
    if user:
        user_configs = get_configs(user["user_id"])

    return {
        "success": True,
        "models": models_list,
        "configs": user_configs,
        "active": active
    }

@app.get("/api/models/default")
def api_get_default_model(user = Depends(get_user_from_header)):
    """获取默认模型"""
    if not user:
        return {"success": False, "error": "未登录"}
    config = get_default_config(user["user_id"])
    if config:
        config["api_key"] = "****" + config["api_key"][-8:]  # 脱敏
    return {"success": True, "config": config}

@app.post("/api/models")
def api_add_model(req: AddModelReq, user = Depends(get_user_from_header)):
    """添加模型配置"""
    if not user:
        return {"success": False, "error": "未登录"}
    return add_config(user["user_id"], req.name, req.api_base, req.api_key, req.model_name, req.provider, req.is_default)

@app.put("/api/models/{config_id}")
def api_update_model(config_id: int, req: UpdateModelReq, user = Depends(get_user_from_header)):
    """更新模型配置"""
    if not user:
        return {"success": False, "error": "未登录"}
    kwargs = {k: v for k, v in {
        "name": req.name, "api_base": req.api_base, "api_key": req.api_key,
        "model_name": req.model_name, "provider": req.provider,
        "is_default": req.is_default, "status": req.status
    }.items() if v is not None}
    return update_config(config_id, user["user_id"], **kwargs)

@app.delete("/api/models/{config_id}")
def api_delete_model(config_id: int, user = Depends(get_user_from_header)):
    """删除模型配置"""
    if not user:
        return {"success": False, "error": "未登录"}
    return delete_config(config_id, user["user_id"])

@app.post("/api/models/{config_id}/default")
def api_set_default_model(config_id: int, user = Depends(get_user_from_header)):
    """设为默认模型"""
    if not user:
        return {"success": False, "error": "未登录"}
    return set_default(config_id, user["user_id"])

@app.post("/api/models/test")
async def api_test_model(
    api_base: str,
    api_key: str,
    model_name: str,
    provider: str = "custom"
):
    """测试模型连接"""
    if not api_key:
        return {"success": False, "error": "API Key不能为空，请填写后再测试"}
    
    config = {
        "api_base": api_base,
        "api_key": api_key,
        "model_name": model_name,
        "provider": provider
    }
    
    # 简单测试：用模型生成一句话
    messages = [{"role": "user", "content": "说一个词：测试"}]
    result = await call_model(config, messages, max_tokens=50)
    
    if result["success"]:
        return {"success": True, "message": "连接成功", "response": result["content"][:100]}
    return {"success": False, "error": result.get("error", "测试失败")}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)